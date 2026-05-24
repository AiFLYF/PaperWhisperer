import concurrent.futures
import hashlib
import ipaddress
import json
import logging
import mimetypes
import os
import re
import secrets
import socket
import ssl
import time
import threading
import unicodedata
import urllib.error
import urllib.parse
import urllib.request
import uuid
import xml.etree.ElementTree as ET
from datetime import datetime

try:
    import certifi
except ImportError:
    certifi = None

import uvicorn
from fastapi import FastAPI, File, Form, Request, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from openai import APIConnectionError, APIStatusError, APITimeoutError, OpenAI
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from env_loader import load_project_env


load_project_env()

logger = logging.getLogger(__name__)

APP_NAME = "PaperWhisperer"
APP_VERSION = os.getenv("PAPERWHISPERER_VERSION", "0.9.0").strip() or "0.9.0"
APP_USER_AGENT = f"{APP_NAME}/{APP_VERSION}"
APP_STARTED_AT = time.time()
STATIC_ICON_CACHE_CONTROL = "public, max-age=86400, immutable"
SECURITY_HEADERS = {
    "X-Content-Type-Options": "nosniff",
    "X-Frame-Options": "DENY",
    "Referrer-Policy": "strict-origin-when-cross-origin",
    "Permissions-Policy": "camera=(), microphone=(), geolocation=()",
}
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
CONTEXT_FOLDER = "context"
MAX_CONTENT_LENGTH = 16 * 1024 * 1024  # 16MB max limit

app = FastAPI(title=APP_NAME, version=APP_VERSION)


@app.middleware("http")
async def add_security_headers(request, call_next):
    started_at = time.perf_counter()
    response = await call_next(request)
    response.headers["X-Process-Time-Ms"] = f"{(time.perf_counter() - started_at) * 1000:.2f}"
    for header, value in SECURITY_HEADERS.items():
        response.headers.setdefault(header, value)
    return response


app.mount("/static", StaticFiles(directory="templates/static"), name="static")
templates = Jinja2Templates(directory="templates")

# Ensure required folders exist
for folder in [UPLOAD_FOLDER, OUTPUT_FOLDER, CONTEXT_FOLDER]:
    os.makedirs(folder, exist_ok=True)


SUPPORTED_EXTENSIONS = (".txt", ".pdf", ".docx", ".pptx")
ALLOWED_EXTENSIONS = set(SUPPORTED_EXTENSIONS)
SUPPORTED_FILE_TYPES_TEXT = ", ".join(SUPPORTED_EXTENSIONS)


def parse_int_env(name, default, min_value=1, max_value=32):
    raw = os.getenv(name, "").strip()
    if not raw:
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    return max(min_value, min(max_value, value))


def parse_bool_env(name, default=False):
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def parse_bool_value(value, default=False):
    if value is None:
        return default
    return str(value).strip().lower() in {"1", "true", "yes", "on"}


def clamp_int_value(value, default, min_value=1, max_value=32):
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = default
    return max(min_value, min(max_value, parsed))


ARXIV_API_URL = "http://export.arxiv.org/api/query"
SEMANTIC_SCHOLAR_SEARCH_URL = "https://api.semanticscholar.org/graph/v1/paper/search"
SEMANTIC_SCHOLAR_TIMEOUT_SECONDS = parse_int_env("SEMANTIC_SCHOLAR_TIMEOUT_SECONDS", default=20, min_value=5, max_value=120)
SEMANTIC_SCHOLAR_MAX_RETRIES = parse_int_env("SEMANTIC_SCHOLAR_MAX_RETRIES", default=3, min_value=1, max_value=6)
PAPER_SEARCH_RESULT_LIMIT = parse_int_env("PAPER_SEARCH_RESULT_LIMIT", default=8, min_value=1, max_value=20)
RECOMMENDATION_RESULT_LIMIT = parse_int_env("RECOMMENDATION_RESULT_LIMIT", default=6, min_value=1, max_value=20)
PAPER_SEARCH_ENABLE_REWRITE = parse_bool_env("PAPER_SEARCH_ENABLE_REWRITE", default=True)
PAPER_SEARCH_REWRITE_MODEL = os.getenv("PAPER_SEARCH_REWRITE_MODEL", "").strip()
SEMANTIC_SCHOLAR_API_KEY = os.getenv("SEMANTIC_SCHOLAR_API_KEY", "").strip()
REMOTE_IMPORT_TIMEOUT_SECONDS = parse_int_env("REMOTE_IMPORT_TIMEOUT_SECONDS", default=30, min_value=5, max_value=180)
SESSION_TTL_SECONDS = parse_int_env("SESSION_TTL_SECONDS", default=24 * 60 * 60, min_value=60, max_value=30 * 24 * 60 * 60)
SESSION_CLEANUP_INTERVAL_SECONDS = parse_int_env("SESSION_CLEANUP_INTERVAL_SECONDS", default=10 * 60, min_value=60, max_value=24 * 60 * 60)
SESSION_PERSIST_FULL_DOCUMENT = parse_bool_env("SESSION_PERSIST_FULL_DOCUMENT", default=False)

MAX_LLM_CONCURRENCY = parse_int_env("OPENAI_MAX_CONCURRENCY", default=5, min_value=1, max_value=32)
LLM_REQUEST_SEMAPHORE = threading.BoundedSemaphore(MAX_LLM_CONCURRENCY)
LAST_SESSION_CLEANUP_AT = 0.0


def resolve_api_key(explicit_key):
    explicit_text = str(explicit_key or "").strip()
    if explicit_text:
        return explicit_text
    return os.getenv("OPENAI_API_KEY", "").strip()


def is_allowed_file(filename):
    return os.path.splitext(filename)[1].lower() in ALLOWED_EXTENSIONS


def secure_filename(filename):
    value = unicodedata.normalize("NFKD", str(filename)).encode("ascii", "ignore").decode("ascii")
    value = value.replace("/", " ").replace("\\", " ")
    value = "_".join(value.split())
    value = re.sub(r"[^A-Za-z0-9_.-]", "", value)
    value = re.sub(r"_+", "_", value)
    value = value.strip("._")
    if os.name == "nt" and value and value.split(".")[0].upper() in {
        "CON", "PRN", "AUX", "NUL", "COM1", "COM2", "COM3", "COM4", "COM5", "COM6", "COM7", "COM8", "COM9",
        "LPT1", "LPT2", "LPT3", "LPT4", "LPT5", "LPT6", "LPT7", "LPT8", "LPT9",
    }:
        value = f"_{value}"
    return value


def sanitize_identifier(raw_value, prefix):
    candidate = secure_filename((raw_value or "").strip())
    return candidate or f"{prefix}_{uuid.uuid4().hex}"


def build_session_id(raw_session_id):
    return sanitize_identifier(raw_session_id, "session")


def build_unique_storage_path(folder, filename):
    root, ext = os.path.splitext(filename)
    return os.path.join(folder, f"{root}_{uuid.uuid4().hex}{ext}")


def build_safe_upload_filename(filename):
    original_ext = os.path.splitext(filename)[1].lower()
    if original_ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Unsupported file type. Please upload one of: {SUPPORTED_FILE_TYPES_TEXT}")

    sanitized = secure_filename(filename)
    if not sanitized or not sanitized.lower().endswith(original_ext):
        return f"file_{uuid.uuid4().hex}{original_ext}"
    return sanitized


def clean_extracted_text(text):
    """Clean extracted text: collapse excessive blank lines and trim whitespace."""
    if not text:
        return ""
    # Collapse 3+ consecutive newlines into 2
    text = re.sub(r'\n{3,}', '\n\n', text)
    # Collapse runs of whitespace (excluding newlines) into single space
    text = re.sub(r'[^\S\n]+', ' ', text)
    # Strip leading/trailing whitespace per line
    lines = [line.strip() for line in text.split('\n')]
    return '\n'.join(lines).strip()


def get_session_file_path(session_id):
    return os.path.join(CONTEXT_FOLDER, f"{session_id}.json")


def now_iso():
    return datetime.now().isoformat(timespec="seconds")


def parse_iso_datetime(value):
    text = str(value or "").strip()
    if not text:
        return None
    try:
        return datetime.fromisoformat(text)
    except ValueError:
        return None


def build_session_expiry(now=None):
    base_time = now or datetime.now()
    return datetime.fromtimestamp(base_time.timestamp() + SESSION_TTL_SECONDS).isoformat(timespec="seconds")


def generate_session_token():
    return secrets.token_urlsafe(24)


def hash_session_token(token):
    return hashlib.sha256(str(token or "").encode("utf-8")).hexdigest()


def build_document_excerpt(content, limit=12000):
    return (content or "")[:limit]


PROMPT_STABLE_PREFIX = """你是 PaperWhisperer 的学术研究助手。目标是帮助研究者快速理解论文、保留证据、识别结构与局限。
始终优先依据用户提供的文档内容，不编造文献细节；文档没有依据时明确说明。
输出默认使用中文；只有检索词、JSON 字段值或代码任务明确要求时才使用英文。"""

PROMPT_FORMULA_RULES = """公式与格式规则：
- 行内公式必须使用 $...$，不要使用 \\( ... \\)。
- 块级公式必须使用 $$...$$，不要使用 \\[ ... \\]。
- 公式内部的下划线和星号不要做 Markdown 转义。
- 不要把公式放进普通代码块。"""

PROMPT_QUALITY_RULES = """回答质量规则：
- 先给结论，再给依据或层级展开。
- 区分文档事实、合理推断和不确定内容。
- 避免空泛套话，优先输出可直接用于阅读、复盘或追问的内容。
- 保持结构清晰，标题层级不要过深。"""

PROMPT_TASK_CONTRACTS = {
    "summary_chunk": """任务契约：从文献片段提取核心观点和高价值引用，保留关键术语、方法、数据集、结论和公式。""",
    "summary_merge": """任务契约：整合多个片段摘要，去重、合并同义观点，并形成一份连贯的整篇论文概要。""",
    "quotes": """任务契约：提取最值得引用的原句或接近原文的关键表述，不要改写成泛泛总结。""",
    "mindmap": """任务契约：识别论文的研究问题、方法、实验、结论与局限，输出文本层级结构。""",
    "mermaid": """任务契约：输出可渲染的 Mermaid 结构图代码，只输出 Mermaid，不解释。""",
    "evaluation": """任务契约：以审稿和读者复盘视角评价论文贡献、优点、局限、历史地位与学习价值。""",
    "research_brief": """任务契约：生成面向研究决策的深度阅读简报，连接贡献、证据、复现风险和后续检索方向。""",
    "qa": """任务契约：基于当前文档和最近问答历史回答用户追问；文档优先级高于历史。""",
    "search_rewrite": """Task contract: rewrite paper-search intent into concise English retrieval queries and return JSON only.""",
}

ANSWER_MODES = {
    "evidence": "证据模式：先给直接结论，再列出文档依据、可推断内容和不确定处。",
    "explain": "讲解模式：用教学方式分步骤解释概念、方法和公式，必要时补充类比。",
    "critique": "评审模式：从贡献、假设、局限、威胁效度和可改进点进行批判性分析。",
    "reproduce": "复现模式：输出复现步骤、关键变量、数据/实验依赖、风险点和检查清单。",
}

READING_QUEUE_LIMIT = 30


def build_stable_system_prompt(task_key):
    task_contract = PROMPT_TASK_CONTRACTS.get(task_key, "任务契约：完成用户指定的学术阅读任务。")
    return "\n\n".join([PROMPT_STABLE_PREFIX, PROMPT_FORMULA_RULES, PROMPT_QUALITY_RULES, task_contract])


def build_prompt_block(tag, content):
    safe_tag = re.sub(r"[^A-Za-z0-9_]", "_", str(tag or "input")).strip("_") or "input"
    return f"<{safe_tag}>\n{str(content or '').strip()}\n</{safe_tag}>"


def build_task_user_prompt(task, input_blocks, constraints="", output_format=""):
    parts = [build_prompt_block("task", task)]
    if constraints:
        parts.append(build_prompt_block("constraints", constraints))
    if output_format:
        parts.append(build_prompt_block("output_format", output_format))
    for tag, content in input_blocks:
        parts.append(build_prompt_block(tag, content))
    parts.append(build_prompt_block("self_check", "提交前确认：没有编造文档外事实；公式格式符合要求；输出结构与任务契约一致。"))
    return "\n\n".join(parts)


def normalize_text_items(values, max_items=8, item_limit=180):
    if not isinstance(values, list):
        return []
    items = []
    seen = set()
    for value in values:
        text = compact_text(value, limit=item_limit)
        if not text or text in seen:
            continue
        items.append(text)
        seen.add(text)
        if len(items) >= max_items:
            break
    return items


def normalize_next_actions(values, max_items=5):
    if not isinstance(values, list):
        return []
    actions = []
    seen = set()
    for value in values:
        if not isinstance(value, dict):
            continue
        label = compact_text(value.get("label"), limit=40)
        prompt = compact_text(value.get("prompt"), limit=240)
        if not label or not prompt or prompt in seen:
            continue
        actions.append({"label": label, "prompt": prompt})
        seen.add(prompt)
        if len(actions) >= max_items:
            break
    return actions


def normalize_answer_mode(value):
    mode = str(value or "").strip().lower()
    return mode if mode in ANSWER_MODES else "evidence"


def normalize_paper_collection(values, max_items=READING_QUEUE_LIMIT):
    if not isinstance(values, list):
        return []
    items = []
    seen = set()
    for value in values:
        if not isinstance(value, dict):
            continue
        title = compact_text(value.get("title"), limit=300)
        if not title:
            continue
        url = str(value.get("url") or "").strip()[:1000]
        pdf_url = str(value.get("pdf_url") or "").strip()[:1000]
        key = re.sub(r"\s+", " ", title).strip().lower() or url or pdf_url
        if key in seen:
            continue
        seen.add(key)
        items.append({
            "source": compact_text(value.get("source"), limit=80),
            "paper_id": compact_text(value.get("paper_id"), limit=160),
            "title": title,
            "abstract": compact_text(value.get("abstract"), limit=1200),
            "authors": normalize_author_list(value.get("authors") or [], limit=8),
            "year": parse_year(value.get("year")),
            "venue": compact_text(value.get("venue"), limit=120),
            "url": url,
            "pdf_url": pdf_url,
            "saved_at": compact_text(value.get("saved_at") or now_iso(), limit=40),
        })
        if len(items) >= max_items:
            break
    return items


def build_analysis_metadata(sections):
    section_labels = {
        "summary": "概览",
        "quotes": "引用片段",
        "mindmap": "文本结构",
        "mermaid": "视觉图谱",
        "evaluation": "批判评价",
        "research_brief": "深度简报",
    }
    section_statuses = {
        name: str((section or {}).get("status") or "empty")
        for name, section in (sections or {}).items()
    }
    completed = [section_labels.get(name, name) for name, status in section_statuses.items() if status == "success"]
    failed = [section_labels.get(name, name) for name, status in section_statuses.items() if status == "failed"]
    disabled = [section_labels.get(name, name) for name, status in section_statuses.items() if status == "disabled"]

    suggested_questions = [
        "这篇论文要解决的核心问题是什么？",
        "它的主要方法和创新点分别是什么？",
        "实验或论证最支持哪些结论？",
        "这篇论文有哪些局限性和后续研究方向？",
    ]
    if section_statuses.get("quotes") == "success":
        suggested_questions.append("哪些原文片段最适合在综述或笔记中引用？")
    if section_statuses.get("evaluation") == "success":
        suggested_questions.append("如果我要复现或扩展这篇论文，应该优先关注什么？")
    if section_statuses.get("research_brief") == "success":
        suggested_questions.append("这篇论文最适合放进哪条研究脉络或综述段落？")

    next_actions = [
        {"label": "追问方法细节", "prompt": "请解释这篇论文的方法流程，并指出每一步解决了什么问题。"},
        {"label": "整理局限性", "prompt": "请基于文档总结这篇论文的局限性，并给出可能的改进方向。"},
        {"label": "生成阅读路线", "prompt": "请把这篇论文拆成适合精读的阅读路线和检查清单。"},
    ]
    if section_statuses.get("research_brief") == "success":
        next_actions.extend([
            {"label": "复现路线", "prompt": "请基于深度简报生成一份复现路线，包括数据、实验变量、依赖和风险点。"},
            {"label": "找后续工作", "prompt": "请提炼 5 个英文检索关键词，用于寻找这篇论文的后续工作或相邻研究。"},
        ])
    if section_statuses.get("evaluation") == "disabled":
        next_actions.append({"label": "手动评价", "prompt": "请基于当前文档补充一份批判性评价。"})
    if failed:
        next_actions.append({"label": "补全失败部分", "prompt": f"请重新生成以下分析部分：{', '.join(failed)}。"})

    return {
        "suggested_questions": suggested_questions[:6],
        "next_actions": next_actions[:5],
        "analysis_status": {
            "quality": "partial" if failed else "complete",
            "completed_sections": completed,
            "failed_sections": failed,
            "disabled_sections": disabled,
            "section_statuses": section_statuses,
        },
    }


def trim_text_for_log(text, limit=2000):
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "\n...[truncated]"


def build_ssl_context():
    if certifi:
        return ssl.create_default_context(cafile=certifi.where())
    return ssl.create_default_context()


def get_retry_delay_seconds(attempt, max_delay=6):
    try:
        attempt_index = max(0, int(attempt))
    except (TypeError, ValueError):
        attempt_index = 0
    try:
        delay_cap = max(1, int(max_delay))
    except (TypeError, ValueError):
        delay_cap = 6
    return min(2 * (attempt_index + 1), delay_cap)


def http_get_json(url, timeout=20, headers=None, retries=1, ssl_context=None):
    request = urllib.request.Request(url, headers=headers or {})
    last_error = None
    for attempt in range(max(1, retries)):
        try:
            with urllib.request.urlopen(request, timeout=timeout, context=ssl_context) as response:
                charset = response.headers.get_content_charset() or "utf-8"
                return json.loads(response.read().decode(charset, errors="ignore"))
        except urllib.error.HTTPError as exc:
            last_error = exc
            if exc.code == 429 and attempt < retries - 1:
                time.sleep(get_retry_delay_seconds(attempt))
                continue
            raise
        except Exception as exc:
            last_error = exc
            if attempt < retries - 1:
                time.sleep(get_retry_delay_seconds(attempt))
                continue
            raise
    raise last_error


def http_get_text(url, timeout=20, headers=None, retries=1, ssl_context=None):
    request = urllib.request.Request(url, headers=headers or {})
    last_error = None
    for attempt in range(max(1, retries)):
        try:
            with urllib.request.urlopen(request, timeout=timeout, context=ssl_context) as response:
                charset = response.headers.get_content_charset() or "utf-8"
                return response.read().decode(charset, errors="ignore")
        except urllib.error.HTTPError as exc:
            last_error = exc
            if exc.code == 429 and attempt < retries - 1:
                time.sleep(get_retry_delay_seconds(attempt))
                continue
            raise
        except Exception as exc:
            last_error = exc
            if attempt < retries - 1:
                time.sleep(get_retry_delay_seconds(attempt))
                continue
            raise
    raise last_error


def compact_text(text, limit=400):
    collapsed = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(collapsed) <= limit:
        return collapsed
    return collapsed[:limit].rstrip() + "..."


def parse_year(value):
    text = str(value or "").strip()
    if not text:
        return ""
    match = re.search(r"(19|20)\d{2}", text)
    return match.group(0) if match else ""


def normalize_author_list(authors, limit=8):
    normalized = []
    for author in authors or []:
        if isinstance(author, str):
            name = author.strip()
        elif isinstance(author, dict):
            name = str(author.get("name") or author.get("author") or "").strip()
        else:
            name = str(getattr(author, "name", "") or "").strip()
        if name:
            normalized.append(name)
        if len(normalized) >= limit:
            break
    return normalized


def normalize_paper_record(source, record):
    if source == "Semantic Scholar":
        open_access_pdf = record.get("openAccessPdf") or {}
        return {
            "source": source,
            "paper_id": str(record.get("paperId") or "").strip(),
            "title": compact_text(record.get("title") or "", limit=300),
            "abstract": compact_text(record.get("abstract") or "", limit=2000),
            "authors": normalize_author_list(record.get("authors") or []),
            "year": parse_year(record.get("year")),
            "venue": compact_text(record.get("venue") or "Semantic Scholar", limit=120),
            "url": str(record.get("url") or "").strip(),
            "pdf_url": str(open_access_pdf.get("url") or "").strip(),
        }

    return {
        "source": source,
        "paper_id": str(record.get("paper_id") or record.get("id") or "").strip(),
        "title": compact_text(record.get("title") or "", limit=300),
        "abstract": compact_text(record.get("abstract") or record.get("summary") or "", limit=2000),
        "authors": normalize_author_list(record.get("authors") or []),
        "year": parse_year(record.get("year") or record.get("published") or ""),
        "venue": compact_text(record.get("venue") or source, limit=120),
        "url": str(record.get("url") or record.get("id") or "").strip(),
        "pdf_url": str(record.get("pdf_url") or "").strip(),
    }


def deduplicate_papers(items):
    deduplicated = []
    seen_titles = set()
    for item in items or []:
        title_key = re.sub(r"\s+", " ", str(item.get("title") or "")).strip().lower()
        if not title_key or title_key in seen_titles:
            continue
        seen_titles.add(title_key)
        deduplicated.append(item)
    return deduplicated


def search_arxiv_papers(query, limit):
    encoded_query = urllib.parse.quote(query)
    url = f"{ARXIV_API_URL}?search_query=all:{encoded_query}&start=0&max_results={limit}"
    feed_text = http_get_text(
        url,
        timeout=SEMANTIC_SCHOLAR_TIMEOUT_SECONDS,
        headers={"User-Agent": APP_USER_AGENT},
        retries=2,
        ssl_context=build_ssl_context(),
    )
    root = ET.fromstring(feed_text)
    namespace = {"atom": "http://www.w3.org/2005/Atom"}
    items = []

    for entry in root.findall("atom:entry", namespace):
        title = compact_text(entry.findtext("atom:title", default="", namespaces=namespace), limit=300)
        summary = compact_text(entry.findtext("atom:summary", default="", namespaces=namespace), limit=2000)
        paper_id = (entry.findtext("atom:id", default="", namespaces=namespace) or "").strip()
        published = (entry.findtext("atom:published", default="", namespaces=namespace) or "").strip()
        authors = [author.findtext("atom:name", default="", namespaces=namespace) for author in entry.findall("atom:author", namespace)]
        pdf_url = ""
        for link in entry.findall("atom:link", namespace):
            if link.attrib.get("title") == "pdf":
                pdf_url = link.attrib.get("href", "").strip()
                break
        items.append(normalize_paper_record("arXiv", {
            "paper_id": paper_id,
            "title": title,
            "abstract": summary,
            "authors": authors,
            "published": published,
            "venue": "arXiv",
            "url": paper_id,
            "pdf_url": pdf_url,
        }))
    return items


def search_semantic_scholar_papers(query, limit):
    params = urllib.parse.urlencode({
        "query": query,
        "limit": limit,
        "fields": "title,abstract,year,venue,url,authors,openAccessPdf,paperId",
    })
    url = f"{SEMANTIC_SCHOLAR_SEARCH_URL}?{params}"
    headers = {"User-Agent": APP_USER_AGENT}
    if SEMANTIC_SCHOLAR_API_KEY:
        headers["x-api-key"] = SEMANTIC_SCHOLAR_API_KEY
    payload = http_get_json(
        url,
        timeout=SEMANTIC_SCHOLAR_TIMEOUT_SECONDS,
        headers=headers,
        retries=SEMANTIC_SCHOLAR_MAX_RETRIES,
    )
    return [normalize_paper_record("Semantic Scholar", item) for item in payload.get("data", [])]


def search_papers(query, limit=None):
    clean_query = compact_text(query, limit=240)
    if not clean_query:
        raise ValueError("Please enter a search query.")

    resolved_limit = clamp_int_value(limit, PAPER_SEARCH_RESULT_LIMIT, min_value=1, max_value=PAPER_SEARCH_RESULT_LIMIT)
    items = []
    errors = []

    for source_name, search_fn in (
        ("Semantic Scholar", search_semantic_scholar_papers),
        ("arXiv", search_arxiv_papers),
    ):
        try:
            items.extend(search_fn(clean_query, resolved_limit))
        except urllib.error.HTTPError as exc:
            logger.warning("%s paper search failed: %s", source_name, exc)
            if exc.code == 429:
                errors.append(f"{source_name}: rate limit reached, please retry in a moment")
            else:
                errors.append(f"{source_name}: HTTP {exc.code}")
        except ssl.SSLCertVerificationError:
            logger.warning("%s paper search SSL verification failed", source_name)
            errors.append(f"{source_name}: SSL certificate verification failed")
        except urllib.error.URLError as exc:
            logger.warning("%s paper search failed: %s", source_name, exc)
            reason = getattr(exc, "reason", exc)
            if isinstance(reason, ssl.SSLCertVerificationError) or "CERTIFICATE_VERIFY_FAILED" in str(reason):
                errors.append(f"{source_name}: SSL certificate verification failed")
            else:
                errors.append(f"{source_name}: {reason}")
        except Exception as exc:
            logger.warning("%s paper search failed: %s", source_name, exc)
            errors.append(f"{source_name}: {exc}")

    return {
        "query": clean_query,
        "items": deduplicate_papers(items)[:resolved_limit],
        "errors": errors,
    }


def build_session_payload(session_id, source_filename, document_content, analysis, session_token):
    generated_at = now_iso()
    stored_document_content = document_content if SESSION_PERSIST_FULL_DOCUMENT else ""
    return {
        "session_id": session_id,
        "source_filename": source_filename,
        "generated_at": generated_at,
        "created_at": generated_at,
        "updated_at": generated_at,
        "expires_at": build_session_expiry(),
        "document_content": stored_document_content,
        "document_excerpt": build_document_excerpt(document_content),
        "qa_history": [],
        "paper_search": {
            "last_query": "",
            "last_results": [],
            "last_recommendation": {},
            "reading_queue": [],
        },
        "session_auth": {
            "token_hash": hash_session_token(session_token),
        },
        "analysis": {
            "summary": analysis.get("summary", ""),
            "quotes": analysis.get("quotes", ""),
            "mindmap": analysis.get("mindmap", ""),
            "mermaid": analysis.get("mermaid", ""),
            "evaluation": analysis.get("evaluation", ""),
            "research_brief": analysis.get("research_brief", ""),
            "sections": analysis.get("sections", {}),
            "char_count": analysis.get("char_count", 0),
            "elapsed_seconds": analysis.get("elapsed_seconds"),
            "output_file": analysis.get("output_file", ""),
            "suggested_questions": analysis.get("suggested_questions", []),
            "next_actions": analysis.get("next_actions", []),
            "analysis_status": analysis.get("analysis_status", {}),
        },
    }


def is_failed_llm_result(value):
    text = (value or "").strip()
    return text.startswith("生成失败，请重试")


def build_section_result(status, content="", error="", retryable=False):
    return {
        "status": status,
        "content": content or "",
        "error": error or "",
        "retryable": bool(retryable),
    }


def build_sse_event(event_name, payload):
    data = json.dumps(payload or {}, ensure_ascii=False)
    return f"event: {event_name}\ndata: {data}\n\n"


def build_sse_headers():
    return {
        "Cache-Control": "no-cache",
        "Connection": "keep-alive",
        "X-Accel-Buffering": "no",
    }


def build_error_payload(message, code="bad_request"):
    return {"error": message, "code": code, "timestamp": now_iso()}


def build_error_response(message, status_code=400, code="bad_request"):
    return JSONResponse(
        content=build_error_payload(message, code=code),
        status_code=status_code,
    )


def describe_remote_http_error(status_code):
    if status_code == 404:
        return "The paper file could not be found at the remote source."
    if status_code == 403:
        return "The remote source denied access to the paper file."
    if status_code == 429:
        return "The remote source rate limited the paper download. Please retry in a moment."
    return f"Remote paper download failed with HTTP {status_code}."


def describe_remote_url_error(reason):
    if isinstance(reason, ssl.SSLCertVerificationError) or "CERTIFICATE_VERIFY_FAILED" in str(reason):
        return "SSL certificate verification failed while downloading the paper file."
    return f"Paper download failed: {reason}"


async def parse_json_object_request(request):
    try:
        raw_body = await request.body()
    except Exception as exc:
        logger.warning("Failed to read JSON request body: %s", exc)
        return None, build_error_response("Unable to read request body.", code="body_read_failed")

    if not raw_body or not raw_body.strip():
        return {}, None

    try:
        data = json.loads(raw_body)
    except json.JSONDecodeError:
        return None, build_error_response("Invalid JSON body.", code="invalid_json")

    if not isinstance(data, dict):
        return None, build_error_response("JSON body must be an object.", code="json_not_object")
    return data, None


def is_retryable_llm_error(message):
    normalized = str(message or "").lower()
    retry_markers = (
        "超时",
        "timeout",
        "连接失败",
        "connection",
        "429",
        "限流",
        "502",
        "503",
        "504",
        "网关",
    )
    return any(marker in normalized for marker in retry_markers)


def describe_llm_status_code(status_code):
    status_messages = {
        400: "AI 服务请求格式错误，请检查模型配置、参数或接口兼容性。",
        401: "AI 服务认证失败，请检查 API Key 是否正确或已过期。",
        403: "AI 服务拒绝访问，当前 API Key 可能无权使用该模型或接口。",
        404: "AI 服务地址或模型不存在，请检查 OPENAI_BASE_URL 和模型名称。",
        408: "AI 服务请求超时，请稍后重试。",
        409: "AI 服务请求冲突，请稍后重试。",
        415: "AI 服务不支持当前请求媒体类型，请检查供应商兼容性。",
        422: "AI 服务无法处理当前请求，请检查输入内容或参数。",
        429: "AI 服务触发限流，请稍后重试或降低并发。",
        500: "AI 服务提供商内部错误，请稍后重试。",
        502: "AI 服务网关异常，请稍后重试。",
        503: "AI 服务暂时不可用，请稍后重试。",
        504: "AI 服务网关超时，请稍后重试。",
    }
    return status_messages.get(status_code, f"AI 服务请求失败，状态码: {status_code}。")


def looks_like_html_response(value):
    text = str(value or "").lstrip().lower()
    html_markers = ("<!doctype html", "<html", "<head", "<body", "<meta ")
    return any(text.startswith(marker) for marker in html_markers)


def extract_message_text(content):
    if content is None:
        return ""
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        text_parts = []
        for item in content:
            if isinstance(item, dict):
                if item.get("type") == "text" and item.get("text"):
                    text_parts.append(str(item["text"]))
            else:
                item_type = getattr(item, "type", None)
                item_text = getattr(item, "text", None)
                if item_type == "text" and item_text:
                    text_parts.append(str(item_text))
        return "\n".join(text_parts).strip()
    return str(content).strip()


def extract_json_object(text):
    raw_text = str(text or "").strip()
    if raw_text.startswith("```"):
        raw_text = re.sub(r"^```(?:json)?\s*", "", raw_text, flags=re.IGNORECASE)
        raw_text = re.sub(r"\s*```$", "", raw_text)
    start = raw_text.find("{")
    end = raw_text.rfind("}")
    if start != -1 and end != -1 and end >= start:
        return raw_text[start:end + 1]
    return raw_text


def remove_file_safely(file_path, description="temporary file"):
    if not file_path or not os.path.exists(file_path):
        return
    try:
        os.remove(file_path)
    except Exception:
        logger.exception("Failed to remove %s: %s", description, file_path)


def close_response_safely(response, description="remote response"):
    if not response:
        return
    try:
        response.close()
    except Exception:
        logger.exception("Failed to close %s", description)


async def close_upload_file_safely(upload_file, description="upload file"):
    if not upload_file:
        return
    try:
        await upload_file.close()
    except Exception:
        logger.exception("Failed to close %s", description)


def atomic_write_json(file_path, payload):
    temp_path = f"{file_path}.{uuid.uuid4().hex}.tmp"
    try:
        with open(temp_path, "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False, indent=2)
            f.flush()
            os.fsync(f.fileno())
        os.replace(temp_path, file_path)
    finally:
        remove_file_safely(temp_path, "stale JSON temp file")


def write_session_payload(session_id, payload):
    now_text = now_iso()
    if isinstance(payload, dict):
        payload["updated_at"] = now_text
        payload["expires_at"] = build_session_expiry()
    atomic_write_json(get_session_file_path(session_id), payload)


def normalize_content_type(value):
    return str(value or "").split(";", 1)[0].strip().lower()


def parse_content_length(headers):
    raw_value = str(headers.get("Content-Length") or "").strip()
    if not raw_value:
        return None
    try:
        value = int(raw_value)
    except ValueError:
        return None
    return value if value >= 0 else None


def validate_remote_content_length(headers, max_bytes):
    content_length = parse_content_length(headers)
    if content_length is None:
        return
    if content_length <= 0:
        raise ValueError("Downloaded file is empty.")
    if content_length > max_bytes:
        raise ValueError(f"Remote file is too large. Limit: {max_bytes // (1024 * 1024)} MB")


def validate_remote_content_type(file_name, content_type):
    normalized = normalize_content_type(content_type)
    if not normalized:
        return

    rejected_types = {"application/json", "application/xhtml+xml", "application/xml", "text/xml"}
    if normalized.startswith("text/html") or normalized in rejected_types:
        raise ValueError("The paper link returned a non-document response instead of a downloadable file.")

    generic_binary_types = {"application/octet-stream", "binary/octet-stream", "application/download", "application/x-download", "application/force-download"}
    allowed_types = {
        ".pdf": {"application/pdf", "application/x-pdf", "application/acrobat", "application/vnd.pdf"},
        ".txt": {"text/plain", "text/markdown"},
        ".docx": {"application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/zip", "application/x-zip-compressed"},
        ".pptx": {"application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/zip", "application/x-zip-compressed"},
    }
    ext = os.path.splitext(file_name)[1].lower()
    if normalized in generic_binary_types or normalized in allowed_types.get(ext, set()):
        return
    if ext == ".txt" and normalized.startswith("text/"):
        return
    raise ValueError(f"Remote content type is not compatible with {ext or 'the selected'} file.")


def get_sample_head(sample):
    head = bytes(sample or b"")[:2048].lstrip()
    if head.startswith(b"\xef\xbb\xbf"):
        head = head[3:].lstrip()
    return head.lower()


def validate_document_file_signature(file_name, sample, source_name):
    if not sample:
        raise ValueError(f"The {source_name} is empty.")

    head = get_sample_head(sample)
    if head.startswith((b"<!doctype html", b"<html")) or b"<html" in head[:512]:
        raise ValueError(f"The {source_name} appears to be an HTML page instead of a document.")

    ext = os.path.splitext(file_name)[1].lower()
    if ext == ".pdf" and b"%PDF-" not in sample[:1024]:
        raise ValueError(f"The {source_name} does not look like a valid PDF.")
    if ext in {".docx", ".pptx"} and not sample.startswith((b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")):
        raise ValueError(f"The {source_name} does not look like a valid {ext[1:].upper()} file.")
    if ext == ".txt" and b"\x00" in sample[:2048]:
        raise ValueError(f"The {source_name} appears to be binary.")


def validate_remote_file_signature(file_name, sample):
    validate_document_file_signature(file_name, sample, "remote file")


def read_remote_file_sample(response, file_name):
    sample = response.read(4096)
    validate_remote_file_signature(file_name, sample)
    return sample


def validate_saved_file_signature(file_path):
    with open(file_path, "rb") as f:
        sample = f.read(4096)
    validate_document_file_signature(os.path.basename(file_path), sample, "uploaded file")


PUBLIC_HOSTNAME_CACHE_TTL_SECONDS = 300
PUBLIC_HOSTNAME_CACHE_MAX_SIZE = 512
PUBLIC_HOSTNAME_CACHE = {}
PUBLIC_HOSTNAME_CACHE_LOCK = threading.Lock()


def is_public_ip_address(value):
    try:
        ip = ipaddress.ip_address(str(value or "").strip())
    except ValueError:
        return False
    return ip.is_global and not (ip.is_loopback or ip.is_link_local or ip.is_multicast or ip.is_private or ip.is_reserved or ip.is_unspecified)


def get_cached_public_hostname_result(normalized):
    now = time.time()
    with PUBLIC_HOSTNAME_CACHE_LOCK:
        cached = PUBLIC_HOSTNAME_CACHE.get(normalized)
        if cached and cached[0] > now:
            return cached[1]
        if cached:
            PUBLIC_HOSTNAME_CACHE.pop(normalized, None)
    return None


def set_cached_public_hostname_result(normalized, result):
    with PUBLIC_HOSTNAME_CACHE_LOCK:
        if len(PUBLIC_HOSTNAME_CACHE) >= PUBLIC_HOSTNAME_CACHE_MAX_SIZE:
            oldest_key = min(PUBLIC_HOSTNAME_CACHE, key=lambda key: PUBLIC_HOSTNAME_CACHE[key][0])
            PUBLIC_HOSTNAME_CACHE.pop(oldest_key, None)
        PUBLIC_HOSTNAME_CACHE[normalized] = (time.time() + PUBLIC_HOSTNAME_CACHE_TTL_SECONDS, bool(result))


def is_ip_literal(value):
    try:
        ipaddress.ip_address(value)
        return True
    except ValueError:
        return False


def resolve_public_hostname(hostname):
    normalized = (hostname or "").strip().strip(".").lower()
    if not normalized or normalized in {"localhost", "localhost.localdomain"} or normalized.endswith(".localhost"):
        return False

    if is_public_ip_address(normalized):
        return True
    if is_ip_literal(normalized):
        return False

    cached_result = get_cached_public_hostname_result(normalized)
    if cached_result is not None:
        return cached_result

    try:
        infos = socket.getaddrinfo(normalized, None, type=socket.SOCK_STREAM)
    except socket.gaierror:
        set_cached_public_hostname_result(normalized, False)
        return False

    resolved_ips = {info[4][0] for info in infos if info and info[4]}
    result = bool(resolved_ips) and all(is_public_ip_address(ip) for ip in resolved_ips)
    set_cached_public_hostname_result(normalized, result)
    return result


def is_public_http_url(raw_url):
    parsed = urllib.parse.urlparse(str(raw_url or "").strip())
    try:
        parsed.port
    except ValueError:
        return False

    if parsed.scheme not in {"http", "https"}:
        return False
    return resolve_public_hostname(parsed.hostname or "")


def looks_like_direct_file_url(raw_url):
    path = urllib.parse.urlparse(str(raw_url or "").strip()).path.lower()
    return any(path.endswith(ext) for ext in SUPPORTED_EXTENSIONS)


def guess_extension_from_content_type(content_type):
    normalized = normalize_content_type(content_type)
    mapping = {
        "application/pdf": ".pdf",
        "text/plain": ".txt",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    }
    return mapping.get(normalized, "")


def extract_filename_from_content_disposition(content_disposition):
    value = str(content_disposition or "")
    match = re.search(r"filename\*=UTF-8''([^;]+)", value, flags=re.IGNORECASE)
    if match:
        return urllib.parse.unquote(match.group(1)).strip('" ')
    match = re.search(r'filename="?([^";]+)"?', value, flags=re.IGNORECASE)
    if match:
        return match.group(1).strip()
    return ""


def build_import_filename(title, source_url, content_disposition, content_type):
    disposition_name = extract_filename_from_content_disposition(content_disposition)
    url_name = os.path.basename(urllib.parse.urlparse(str(source_url or "")).path)
    title_slug = secure_filename(title or "")
    candidate_name = disposition_name or url_name or title_slug or "imported_paper"
    candidate_root, candidate_ext = os.path.splitext(candidate_name)
    inferred_ext = candidate_ext.lower() if candidate_ext.lower() in ALLOWED_EXTENSIONS else ""
    if not inferred_ext:
        inferred_ext = guess_extension_from_content_type(content_type)
    if not inferred_ext:
        inferred_ext = ".pdf"
    safe_root = secure_filename(candidate_root) or title_slug or "imported_paper"
    return build_safe_upload_filename(f"{safe_root}{inferred_ext}")


def iter_downloadable_paper_urls(pdf_url, url):
    seen = set()
    for candidate, require_direct_file in ((pdf_url, False), (url, True)):
        normalized_candidate = str(candidate or "").strip()
        if not normalized_candidate or normalized_candidate in seen:
            continue
        seen.add(normalized_candidate)
        yield normalized_candidate, require_direct_file


def stream_remote_paper(title, pdf_url, url):
    candidate_urls = list(iter_downloadable_paper_urls(pdf_url, url))
    if not candidate_urls:
        raise ValueError("No downloadable paper link found for this result.")

    ssl_context = build_ssl_context()
    last_error = None

    for source_url, require_direct_file in candidate_urls:
        if not is_public_http_url(source_url):
            last_error = ValueError("Only public http/https paper URLs are allowed.")
            continue
        if require_direct_file and not looks_like_direct_file_url(source_url):
            last_error = ValueError("This result does not provide a direct downloadable file. Please open it manually and upload the paper file.")
            continue

        response = None
        request = urllib.request.Request(source_url, headers={"User-Agent": APP_USER_AGENT})
        try:
            response = urllib.request.urlopen(request, timeout=REMOTE_IMPORT_TIMEOUT_SECONDS, context=ssl_context)
            final_url = response.geturl() or source_url
            if not is_public_http_url(final_url):
                raise ValueError("The paper link redirected to a non-public URL.")

            content_type = normalize_content_type(response.headers.get("Content-Type"))
            validate_remote_content_length(response.headers, MAX_CONTENT_LENGTH)
            file_name = build_import_filename(
                title=title,
                source_url=final_url,
                content_disposition=response.headers.get("Content-Disposition"),
                content_type=content_type,
            )
            if not is_allowed_file(file_name):
                raise ValueError(f"Unsupported remote file type. Please use one of: {SUPPORTED_FILE_TYPES_TEXT}")
            validate_remote_content_type(file_name, content_type)
            initial_chunk = read_remote_file_sample(response, file_name)

            return response, file_name, content_type, initial_chunk
        except Exception as exc:
            close_response_safely(response, "failed remote paper response")
            last_error = exc

    if last_error:
        raise last_error
    raise ValueError("Paper import failed.")


async def save_upload_file(upload_file, destination_path, max_bytes):
    total_bytes = 0
    try:
        with open(destination_path, "wb") as f:
            while True:
                chunk = await upload_file.read(1024 * 64)
                if not chunk:
                    break
                total_bytes += len(chunk)
                if total_bytes > max_bytes:
                    raise ValueError(f"Uploaded file is too large. Limit: {max_bytes // (1024 * 1024)} MB")
                f.write(chunk)
        if total_bytes <= 0:
            raise ValueError("Uploaded file is empty.")
        validate_saved_file_signature(destination_path)
        return total_bytes
    except Exception:
        remove_file_safely(destination_path, "failed upload file")
        raise


def iter_remote_file_chunks(response, max_bytes, initial_chunk=b""):
    total_bytes = 0
    try:
        if initial_chunk:
            total_bytes += len(initial_chunk)
            if total_bytes > max_bytes:
                raise ValueError(f"Remote file is too large. Limit: {max_bytes // (1024 * 1024)} MB")
            yield initial_chunk

        while True:
            chunk = response.read(1024 * 64)
            if not chunk:
                break
            total_bytes += len(chunk)
            if total_bytes > max_bytes:
                raise ValueError(f"Remote file is too large. Limit: {max_bytes // (1024 * 1024)} MB")
            yield chunk
        if total_bytes <= 0:
            raise ValueError("Downloaded file is empty.")
    finally:
        close_response_safely(response, "remote file stream response")


def cleanup_expired_sessions(force=False):
    global LAST_SESSION_CLEANUP_AT
    now_ts = time.time()
    if not force and now_ts - LAST_SESSION_CLEANUP_AT < SESSION_CLEANUP_INTERVAL_SECONDS:
        return
    LAST_SESSION_CLEANUP_AT = now_ts
    for name in os.listdir(CONTEXT_FOLDER):
        if not name.endswith(".json"):
            continue
        file_path = os.path.join(CONTEXT_FOLDER, name)
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                payload = json.load(f)
            if not isinstance(payload, dict):
                raise ValueError("session payload must be an object")
            expires_at = parse_iso_datetime(payload.get("expires_at"))
            if expires_at and expires_at.timestamp() < now_ts:
                remove_file_safely(file_path, "expired session cleanup file")
                logger.info("Removed expired session file: %s", file_path)
        except (OSError, json.JSONDecodeError, TypeError, ValueError) as exc:
            remove_file_safely(file_path, "unreadable session cleanup file")
            logger.warning("Removed unreadable session file %s: %s", file_path, exc)
            continue


def validate_session_token(session_payload, session_token):
    expected_hash = str((session_payload.get("session_auth") or {}).get("token_hash") or "")
    provided_hash = hash_session_token(session_token)
    return bool(expected_hash and session_token and secrets.compare_digest(expected_hash, provided_hash))


def get_session_document_content(session_payload):
    content = str(session_payload.get("document_content") or "")
    if content:
        return content
    return str(session_payload.get("document_excerpt") or "")


def load_validated_session(raw_session_id, session_token, require_token=True):
    if not raw_session_id:
        raise ValueError("session_id is required.")
    safe_session_id = build_session_id(raw_session_id)
    session_payload = load_session_payload(safe_session_id)
    if not session_payload:
        raise ValueError("Session expired or context not found. Please upload and analyze the file again.")
    if require_token and not validate_session_token(session_payload, session_token):
        raise PermissionError("Invalid or missing session token. Please analyze the document again.")
    return safe_session_id, session_payload


def download_remote_paper(title, pdf_url, url):
    response = None
    temp_path = None
    try:
        response, file_name, _content_type, initial_chunk = stream_remote_paper(title=title, pdf_url=pdf_url, url=url)
        temp_path = build_unique_storage_path(UPLOAD_FOLDER, file_name)
        with open(temp_path, "wb") as f:
            for chunk in iter_remote_file_chunks(response, MAX_CONTENT_LENGTH, initial_chunk):
                f.write(chunk)

        return temp_path, file_name
    except Exception:
        remove_file_safely(temp_path, "failed remote download file")
        raise
    finally:
        close_response_safely(response, "remote paper download response")


def finalize_analysis_result(result, whisperer, original_filename, generate_evaluation_bool, session_id, generate_research_brief_bool=True):
    safe_session_id = build_session_id(session_id)
    session_token = generate_session_token()
    base_name = os.path.splitext(original_filename)[0]
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = os.path.join(OUTPUT_FOLDER, f"{base_name}_analysis_{timestamp}.md")

    md_content = f"""# PaperWhisperer 分析报告

> 生成时间: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
> 源文件: {original_filename}
> 耗时: {result.get('elapsed_seconds', 'N/A')}s

---

## AI 摘要

{result.get('summary', '')}

---

## 引用片段

{result.get('quotes', '')}

---

## 思维导图

{result.get('mindmap', '')}

---

"""

    if generate_evaluation_bool:
        md_content += f"## 论文评价\n\n{result.get('evaluation', '')}\n\n---\n"

    if generate_research_brief_bool:
        md_content += f"## 深度阅读简报\n\n{result.get('research_brief', '')}\n\n---\n"

    md_content += f"## 元信息\n\n- 版本: {whisperer.version}\n- 字符数: {result['char_count']}\n"

    with open(output_file, "w", encoding="utf-8") as f:
        f.write(md_content)

    result["output_file"] = output_file
    result["session_id"] = safe_session_id
    result["session_token"] = session_token
    result["source_filename"] = original_filename

    session_payload = build_session_payload(
        session_id=safe_session_id,
        source_filename=original_filename,
        document_content=whisperer.document_content,
        analysis=result,
        session_token=session_token,
    )
    write_session_payload(safe_session_id, session_payload)
    return result


def analyze_saved_file(file_path, original_filename, api_key, generate_mermaid_bool, generate_evaluation_bool, session_id, generate_research_brief_bool=True):
    cleanup_expired_sessions()
    resolved_api_key = resolve_api_key(api_key)
    if not resolved_api_key:
        raise ValueError("API key is required. Provide api_key or set OPENAI_API_KEY.")

    whisperer = PaperWhisperer(resolved_api_key)
    result = whisperer.analyze(file_path, generate_mermaid_bool, generate_evaluation_bool, generate_research_brief_bool)
    return finalize_analysis_result(
        result=result,
        whisperer=whisperer,
        original_filename=original_filename,
        generate_evaluation_bool=generate_evaluation_bool,
        session_id=session_id,
        generate_research_brief_bool=generate_research_brief_bool,
    )


def load_session_payload(session_id):
    session_file = get_session_file_path(session_id)
    if not os.path.exists(session_file):
        return None
    try:
        with open(session_file, "r", encoding="utf-8") as f:
            payload = json.load(f)
    except json.JSONDecodeError as exc:
        remove_file_safely(session_file, "corrupt session file")
        logger.warning("Removed corrupt session file %s: %s", session_file, exc)
        return None
    except OSError as exc:
        logger.warning("Unable to read session file %s: %s", session_file, exc)
        return None

    if not isinstance(payload, dict):
        remove_file_safely(session_file, "non-object session file")
        logger.warning("Removed non-object session file %s", session_file)
        return None

    expires_at = parse_iso_datetime(payload.get("expires_at"))
    if expires_at and expires_at.timestamp() < time.time():
        remove_file_safely(session_file, "expired session file")
        return None
    if not expires_at:
        payload["expires_at"] = build_session_expiry()

    document_content = str(payload.get("document_content") or "")
    document_excerpt = str(payload.get("document_excerpt") or build_document_excerpt(document_content))
    qa_history = payload.get("qa_history")
    analysis = payload.get("analysis")
    paper_search = payload.get("paper_search")
    session_auth = payload.get("session_auth")

    payload["document_content"] = document_content
    payload["document_excerpt"] = document_excerpt
    payload["qa_history"] = qa_history if isinstance(qa_history, list) else []
    payload["analysis"] = analysis if isinstance(analysis, dict) else {}
    payload["paper_search"] = paper_search if isinstance(paper_search, dict) else {}
    payload["session_auth"] = session_auth if isinstance(session_auth, dict) else {}
    payload["source_filename"] = str(payload.get("source_filename") or "")
    payload["session_id"] = str(payload.get("session_id") or session_id)
    payload.setdefault("generated_at", now_iso())
    payload.setdefault("created_at", payload.get("generated_at") or now_iso())
    payload.setdefault("updated_at", payload.get("generated_at") or now_iso())
    payload.setdefault("expires_at", build_session_expiry())
    payload["paper_search"].setdefault("last_query", "")
    if not isinstance(payload["paper_search"].get("last_results"), list):
        payload["paper_search"]["last_results"] = []
    if not isinstance(payload["paper_search"].get("last_recommendation"), dict):
        payload["paper_search"]["last_recommendation"] = {}
    payload["paper_search"]["reading_queue"] = normalize_paper_collection(
        payload["paper_search"].get("reading_queue"),
        max_items=READING_QUEUE_LIMIT,
    )
    payload["session_auth"].setdefault("token_hash", "")
    if not isinstance(payload["analysis"].get("sections"), dict):
        payload["analysis"]["sections"] = {}
    payload["analysis"]["suggested_questions"] = normalize_text_items(
        payload["analysis"].get("suggested_questions"),
        max_items=6,
        item_limit=180,
    )
    payload["analysis"]["next_actions"] = normalize_next_actions(payload["analysis"].get("next_actions"), max_items=5)
    if not isinstance(payload["analysis"].get("analysis_status"), dict):
        payload["analysis"]["analysis_status"] = {}
    return payload


class TextChunker:
    """文本分块器，用于处理超长文本"""
    def __init__(self, chunk_size=4000, overlap=200):
        if chunk_size <= 0:
            raise ValueError("chunk_size must be > 0")
        if overlap < 0:
            raise ValueError("overlap must be >= 0")
        if overlap >= chunk_size:
            raise ValueError("overlap must be smaller than chunk_size")
        self.chunk_size = chunk_size
        self.overlap = overlap

    def chunk_text(self, text):
        if len(text) <= self.chunk_size:
            return [text]

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = start + self.chunk_size
            chunk = text[start:end]

            # Try sentence boundaries: Chinese period, English period+space, newline
            if end < text_length:
                best_pos = -1
                for sep in ('。', '. ', '\n'):
                    pos = chunk.rfind(sep)
                    if pos > best_pos:
                        best_pos = pos

                if best_pos > self.chunk_size // 2:
                    chunk = chunk[:best_pos + 1]
                    end = start + best_pos + 1

            if chunk.strip():
                chunks.append(chunk.strip())
            start = end - self.overlap

        return chunks


class DocumentLoader:
    """文档加载器，支持 TXT/PDF/DOCX/PPTX"""
    @staticmethod
    def load_txt(file_path):
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    @staticmethod
    def load_pdf(file_path):
        try:
            reader = PdfReader(file_path)
            text = "\n\n".join([page.extract_text() or "" for page in reader.pages])
            return text.strip()
        except Exception as e:
            raise ValueError(f"PDF 读取失败: {str(e)}")

    @staticmethod
    def load_docx(file_path):
        try:
            doc = Document(file_path)
            parts = []

            # Extract paragraphs
            for p in doc.paragraphs:
                text = p.text.strip()
                if text:
                    parts.append(text)

            # Extract tables
            for table_index, table in enumerate(doc.tables, start=1):
                rows_text = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append(" | ".join(cells))
                if rows_text:
                    parts.append(f"[Table {table_index}]\n" + "\n".join(rows_text))

            return "\n\n".join(parts).strip()
        except Exception as e:
            raise ValueError(f"DOCX 读取失败: {str(e)}")

    @staticmethod
    def load_pptx(file_path):
        try:
            presentation = Presentation(file_path)
            slides_text = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                shape_texts = []
                for shape in slide.shapes:
                    text = getattr(shape, 'text', '')
                    if text and text.strip():
                        shape_texts.append(text.strip())

                # Extract slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        shape_texts.append(f"[Notes] {notes}")

                if shape_texts:
                    slides_text.append(f"[Slide {slide_index}]\n" + "\n".join(shape_texts))
            return "\n\n".join(slides_text).strip()
        except Exception as e:
            raise ValueError(f"PPTX 读取失败: {str(e)}")

    @staticmethod
    def load(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        loaders = {
            '.txt': DocumentLoader.load_txt,
            '.pdf': DocumentLoader.load_pdf,
            '.docx': DocumentLoader.load_docx,
            '.pptx': DocumentLoader.load_pptx,
        }
        loader = loaders.get(ext)
        if not loader:
            raise ValueError(f"不支持的文件格式: {ext} (支持: {SUPPORTED_FILE_TYPES_TEXT})")
        raw_text = loader(file_path)
        return clean_extracted_text(raw_text)


class PaperWhisperer:
    """文献分析核心类"""
    def __init__(self, api_key):
        self.name = APP_NAME
        self.version = APP_VERSION
        self.api_key = resolve_api_key(api_key)
        self.base_url = os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1").strip()
        self.model = os.getenv("OPENAI_MODEL", "gpt-4o-mini").strip()
        self.search_rewrite_model = PAPER_SEARCH_REWRITE_MODEL or self.model
        self.request_timeout = parse_int_env("OPENAI_REQUEST_TIMEOUT_SECONDS", default=60, min_value=5, max_value=600)
        self.max_retries = parse_int_env("OPENAI_MAX_RETRIES", default=3, min_value=1, max_value=10)
        self.chunker = TextChunker(4000, 200)
        self.max_concurrency = MAX_LLM_CONCURRENCY
        self.summary_chunk_workers = min(3, self.max_concurrency)
        self.analysis_workers = min(5, self.max_concurrency)
        self.document_content = ""

        self.client = OpenAI(
            api_key=self.api_key,
            base_url=self.base_url
        ) if self.api_key else None

    def _call_llm(self, system_prompt, user_prompt, max_retries=None, model=None):
        if not self.client:
            raise ValueError("API key is required. Provide it in request body or set OPENAI_API_KEY.")

        retries = self.max_retries if max_retries is None else max_retries

        for attempt in range(retries):
            try:
                with LLM_REQUEST_SEMAPHORE:
                    raw_response = self.client.chat.completions.with_raw_response.create(
                        model=(model or self.model),
                        messages=[
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_prompt}
                        ],
                        temperature=0.7,
                        max_tokens=4000,
                        timeout=self.request_timeout
                    )

                status_code = getattr(raw_response, "status_code", None)
                response = raw_response.parse()

                if status_code is None:
                    raise ValueError("AI 服务未返回可识别的 HTTP 状态码。")
                if not (200 <= status_code < 300):
                    raise ValueError(describe_llm_status_code(status_code))

                if isinstance(response, str):
                    if looks_like_html_response(response):
                        raise ValueError("AI 服务返回了网页内容而不是模型结果。通常是 API Key 缺失、无效，或 OPENAI_BASE_URL 指向了网页地址。")
                    raise ValueError("AI 服务返回了字符串而不是标准响应对象，请检查供应商接口兼容性。")

                choices = getattr(response, "choices", None)
                if not choices:
                    if looks_like_html_response(response):
                        raise ValueError("AI 服务返回了网页内容而不是模型结果。通常是 API Key 缺失、无效，或 OPENAI_BASE_URL 指向了网页地址。")
                    raise ValueError("AI 服务返回成功，但响应中缺少 choices 字段。")

                message = getattr(choices[0], "message", None)
                content = getattr(message, "content", None) if message else None
                text_content = extract_message_text(content)
                if looks_like_html_response(text_content):
                    raise ValueError("AI 服务返回了网页内容而不是模型结果。通常是 API Key 缺失、无效，或 OPENAI_BASE_URL 指向了网页地址。")
                if not text_content:
                    raise ValueError("AI 服务返回内容为空")
                return text_content
            except APIStatusError as e:
                message = describe_llm_status_code(e.status_code)
            except APITimeoutError:
                message = "AI 服务请求超时，请稍后重试。"
            except APIConnectionError:
                message = "AI 服务连接失败，请检查网络、API 地址或供应商服务状态。"
            except Exception as e:
                message = str(e)

            if attempt < retries - 1:
                time.sleep(get_retry_delay_seconds(attempt, max_delay=8))
            else:
                logger.error(message)
                raise RuntimeError(message)

    def _stream_llm(self, system_prompt, user_prompt, max_retries=None, model=None):
        if not self.client:
            raise ValueError("API key is required. Provide it in request body or set OPENAI_API_KEY.")

        retries = self.max_retries if max_retries is None else max_retries

        for attempt in range(retries):
            try:
                with LLM_REQUEST_SEMAPHORE:
                    with self.client.chat.completions.with_streaming_response.create(
                        model=(model or self.model),
                        messages=[
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_prompt}
                        ],
                        temperature=0.7,
                        max_tokens=4000,
                        timeout=self.request_timeout,
                        stream=True,
                    ) as response:
                        status_code = getattr(response, "status_code", None)
                        if status_code is None:
                            raise ValueError("AI 服务未返回可识别的 HTTP 状态码。")
                        if not (200 <= status_code < 300):
                            raise ValueError(describe_llm_status_code(status_code))

                        saw_text = False
                        for chunk in response.iter_lines():
                            if not chunk:
                                continue
                            decoded = chunk.decode("utf-8") if isinstance(chunk, bytes) else str(chunk)
                            line = decoded.strip()
                            if not line.startswith("data:"):
                                continue
                            data = line[5:].strip()
                            if not data or data == "[DONE]":
                                continue
                            try:
                                payload = json.loads(data)
                            except json.JSONDecodeError:
                                continue
                            choices = payload.get("choices") or []
                            if not choices:
                                continue
                            delta = choices[0].get("delta") or {}
                            text = delta.get("content")
                            if not text:
                                continue
                            saw_text = True
                            yield str(text)

                        if not saw_text:
                            raise ValueError("AI 服务返回内容为空")
                        return
            except APIStatusError as e:
                message = describe_llm_status_code(e.status_code)
            except APITimeoutError:
                message = "AI 服务请求超时，请稍后重试。"
            except APIConnectionError:
                message = "AI 服务连接失败，请检查网络、API 地址或供应商服务状态。"
            except Exception as e:
                message = str(e)

            if attempt < retries - 1:
                time.sleep(get_retry_delay_seconds(attempt, max_delay=8))
            else:
                logger.error(message)
                raise RuntimeError(message)

    def _get_worker_count(self, task_count, configured_workers):
        return max(1, min(task_count, configured_workers))

    def _generate_summary_chunk(self, content):
        system_prompt = build_stable_system_prompt("summary_chunk")
        user_prompt = build_task_user_prompt(
            task="从文献片段中提取 3-5 个核心观点，并找出 2-3 个最值得引用的片段。",
            constraints="每个核心观点必须是一句话；引用片段尽量保留原文措辞；不要补充文档外背景。",
            output_format="""## 核心观点
1. [观点1]
2. [观点2]
3. [观点3]

## 引用片段
- "[引用1]"
- "[引用2]""",
            input_blocks=[("document_excerpt", content)],
        )
        return self._call_llm(system_prompt, user_prompt)

    def _merge_summaries(self, summaries):
        if not summaries:
            return None
        if len(summaries) == 1:
            return summaries[0]

        combined = "\n\n--- 章节 ---\n\n".join(summaries)
        system_prompt = build_stable_system_prompt("summary_merge")
        user_prompt = build_task_user_prompt(
            task="把长文献不同片段的摘要整合成一份完整、连贯、去重后的论文概要。",
            constraints="合并重复观点；保留关键方法、贡献、实验结论和引用片段；不要添加片段中没有的信息。",
            output_format="""## 核心观点
1. [整合后的观点1]
2. [整合后的观点2]

## 引用片段
- "[整合后的引用1]"
- "[整合后的引用2]""",
            input_blocks=[("chunk_summaries", combined)],
        )
        return self._call_llm(system_prompt, user_prompt)

    def generate_summary(self, content):
        chunks = self.chunker.chunk_text(content)

        if len(chunks) == 1:
            return self._generate_summary_chunk(content)

        worker_count = self._get_worker_count(len(chunks), self.summary_chunk_workers)
        if worker_count == 1:
            chunk_summaries = [summary for summary in map(self._generate_summary_chunk, chunks) if summary]
        else:
            with concurrent.futures.ThreadPoolExecutor(max_workers=worker_count) as executor:
                chunk_summaries = list(filter(None, executor.map(self._generate_summary_chunk, chunks)))

        if len(chunk_summaries) > 1:
            return self._merge_summaries(chunk_summaries)
        return chunk_summaries[0] if chunk_summaries else "无法生成摘要"

    def extract_quotes(self, content):
        system_prompt = build_stable_system_prompt("quotes")
        user_prompt = build_task_user_prompt(
            task="从文献中提取 3-5 个最值得引用的原句、定义、结论或核心观点。",
            constraints="优先选择能支撑论文主张的方法、发现或结论；尽量保留原文措辞；不要把普通摘要改写成引用。",
            output_format="""## 引用片段
1. "[原句1]"
2. "[原句2]"
3. "[原句3]""",
            input_blocks=[("document_excerpt", build_document_excerpt(content, limit=15000))],
        )
        return self._call_llm(system_prompt, user_prompt)

    def generate_mindmap(self, content):
        system_prompt = build_stable_system_prompt("mindmap")
        user_prompt = build_task_user_prompt(
            task="为文献生成文本格式的研究结构图，帮助用户快速定位论文逻辑。",
            constraints="覆盖研究问题、核心方法、实验或论证、关键结论、局限性；层级控制在 3-4 层；节点短句化。",
            output_format="""## 思维导图
论文主题
├── 研究问题
├── 方法框架
│   ├── [关键模块]
│   └── [关键模块]
├── 证据与实验
└── 结论与局限""",
            input_blocks=[("document_excerpt", build_document_excerpt(content, limit=10000))],
        )
        return self._call_llm(system_prompt, user_prompt)

    def generate_mermaid_mindmap(self, content):
        system_prompt = build_stable_system_prompt("mermaid")
        user_prompt = build_task_user_prompt(
            task="生成可渲染的 Mermaid 论文结构图代码。",
            constraints="必须以 graph TD 或 graph LR 开头；节点 ID 只能包含字母、数字和下划线；节点文本使用方括号；不超过 20 个节点；节点文本不要包含复杂 LaTeX 公式；不要输出 Markdown 代码围栏。",
            output_format="""graph TD
    A[论文标题]
    A --> B[研究问题]
    A --> C[方法]
    A --> D[实验]
    A --> E[结论]""",
            input_blocks=[("document_excerpt", build_document_excerpt(content, limit=4000))],
        )
        result = self._call_llm(system_prompt, user_prompt)
        if result:
            # 寻找真正的 Mermaid 代码起始行，过滤掉大模型输出的开头废话
            lines = result.strip().split('\n')
            start_idx = -1
            valid_prefixes = ("graph ", "mindmap", "flowchart ", "pie", "sequenceDiagram", "stateDiagram", "classDiagram")

            for i, line in enumerate(lines):
                if any(line.strip().startswith(prefix) for prefix in valid_prefixes):
                    start_idx = i
                    break

            if start_idx != -1:
                result = '\n'.join(lines[start_idx:]).strip()
            else:
                # 极端情况下：正则回退提取
                match = re.search(r'```(?:mermaid)?\s*\n(.*?)\n```', result, re.DOTALL | re.IGNORECASE)
                if match:
                    result = match.group(1).strip()
                else:
                    # 默认添加graph TD前缀
                    result = "graph TD\n" + result

            # 确保代码有效
            if not any(result.startswith(prefix) for prefix in valid_prefixes):
                result = "graph TD\n" + result
            return result
        return None

    def generate_evaluation(self, content):
        system_prompt = build_stable_system_prompt("evaluation")
        user_prompt = build_task_user_prompt(
            task="对文献做总结性评价，兼顾审稿视角、读者复盘和后续研究启发。",
            constraints="贡献和局限必须能从文档内容推出；历史地位不确定时说明不确定；避免泛泛而谈。",
            output_format="""## 论文评价

### 主要贡献
[评价内容]

### 历史地位
[评价内容]

### 主要优点
- 优点1
- 优点2

### 局限性
- 局限性1
- 局限性2

### 值得学习的地方
- 学习点1
- 学习点2""",
            input_blocks=[("document_excerpt", build_document_excerpt(content, limit=15000))],
        )
        return self._call_llm(system_prompt, user_prompt)

    def generate_research_brief(self, content):
        system_prompt = build_stable_system_prompt("research_brief")
        user_prompt = build_task_user_prompt(
            task="生成一份可直接用于组会、文献综述和后续检索决策的深度阅读简报。",
            constraints="所有结论必须能从文档内容推出；证据不足时标注不确定；推荐检索词用英文短语；避免泛泛背景介绍。",
            output_format="""## 深度阅读简报

### 一句话定位
[这篇论文解决什么问题、适合放在哪条研究脉络]

### 核心贡献与适用场景
- 贡献1：对应证据或章节线索
- 贡献2：对应证据或章节线索

### 证据-结论链
| 结论 | 文档依据 | 可信度 |
| --- | --- | --- |
| [结论] | [依据] | 高/中/低 |

### 复现检查清单
- 数据与输入要求
- 方法/模型关键变量
- 实验或评估指标
- 潜在失败点

### 局限与后续问题
- 局限1
- 后续问题1

### 推荐检索关键词
- keyword phrase 1
- keyword phrase 2
- keyword phrase 3""",
            input_blocks=[("document_excerpt", build_document_excerpt(content, limit=18000))],
        )
        return self._call_llm(system_prompt, user_prompt)

    def _build_answer_prompts(self, question, history=None, answer_mode="evidence"):
        if not self.document_content:
            raise ValueError("没有文档内容，请先上传文档进行分析。")

        history = history or []
        document_budget = 12000
        total_history_budget = 8000
        per_turn_budget = 2400
        document_window = (self.document_content or "")[:document_budget]

        history_sections = []
        used_history_chars = 0
        for turn in reversed(history):
            question_text = trim_text_for_log(turn.get("question", ""), limit=400)
            answer_text = trim_text_for_log(turn.get("answer", ""), limit=800)
            if not question_text and not answer_text:
                continue

            section = f"Q: {question_text}\nA: {answer_text}"
            section_length = len(section)
            if section_length > per_turn_budget:
                section = section[:per_turn_budget].rstrip() + "\n...[truncated]"
                section_length = len(section)

            if used_history_chars + section_length > total_history_budget:
                break

            history_sections.append(section)
            used_history_chars += section_length

        history_sections.reverse()
        history_block = "\n\n---\n\n".join(history_sections)

        answer_mode = normalize_answer_mode(answer_mode)
        system_prompt = build_stable_system_prompt("qa")
        constraints = (
            "优先依据 document_excerpt；history 只用于理解追问上下文；"
            "如果文档没有答案，明确说明缺少依据；回答要简洁但保留关键证据。"
            f"\n{ANSWER_MODES[answer_mode]}"
        )
        input_blocks = [("document_excerpt", document_window)]
        if history_block:
            input_blocks.append(("recent_qa_history", history_block))
        input_blocks.append(("user_question", question))
        user_prompt = build_task_user_prompt(
            task="回答用户关于当前文档的追问。",
            constraints=constraints,
            output_format="先给直接答案；必要时用要点列出依据、公式或不确定处。",
            input_blocks=input_blocks,
        )
        return system_prompt, user_prompt

    def answer_question(self, question, history=None, answer_mode="evidence"):
        system_prompt, user_prompt = self._build_answer_prompts(question, history=history, answer_mode=answer_mode)
        return self._call_llm(system_prompt, user_prompt)

    def stream_answer_question(self, question, history=None, answer_mode="evidence"):
        system_prompt, user_prompt = self._build_answer_prompts(question, history=history, answer_mode=answer_mode)
        full_answer = []
        for chunk in self._stream_llm(system_prompt, user_prompt):
            full_answer.append(chunk)
            yield chunk
        return "".join(full_answer)

    def rewrite_search_query(self, query, context_text=""):
        clean_query = compact_text(query, limit=240)
        if not clean_query:
            raise ValueError("Please enter a search query.")

        context_excerpt = build_document_excerpt(context_text or "", limit=4000)
        system_prompt = build_stable_system_prompt("search_rewrite")
        user_prompt = build_task_user_prompt(
            task="Rewrite the paper search request into a concise English query for Semantic Scholar and arXiv.",
            constraints="""- rewritten_query must be concise English.
- Preserve specific-paper intent; if a nickname or shorthand points to a known paper, prefer the canonical title.
- Do not broaden a specific-paper query into a vague family query.
- Only expand when the intent is ambiguous.
- topics must be short.
- why must be Chinese.
- Return JSON only, without markdown fences.""",
            output_format='''{
  "original_query": "original user query",
  "rewritten_query": "better english academic query",
  "topics": ["topic 1", "topic 2", "topic 3"],
  "why": "brief reason in Chinese"
}''',
            input_blocks=[
                ("user_query", clean_query),
                ("optional_context", context_excerpt),
            ],
        )
        raw_response = self._call_llm(system_prompt, user_prompt, model=self.search_rewrite_model)
        try:
            rewrite_meta = json.loads(extract_json_object(raw_response))
        except json.JSONDecodeError as exc:
            raise ValueError(f"Search query rewriting failed: {exc}") from exc

        rewritten_query = compact_text(rewrite_meta.get("rewritten_query") or "", limit=240)
        if not rewritten_query:
            raise ValueError("Search query rewriting returned an empty rewritten query.")

        return {
            "original_query": clean_query,
            "rewritten_query": rewritten_query,
            "topics": rewrite_meta.get("topics") or [],
            "reason": str(rewrite_meta.get("why") or "").strip(),
            "model": self.search_rewrite_model,
        }

    def recommend_papers(self, content, limit=None):
        excerpt = build_document_excerpt(content, limit=12000)
        if not excerpt:
            raise ValueError("Current session does not contain document content.")

        resolved_limit = clamp_int_value(limit, RECOMMENDATION_RESULT_LIMIT, min_value=1, max_value=RECOMMENDATION_RESULT_LIMIT)
        rewrite_meta = self.rewrite_search_query(
            query="Find closely related follow-up papers for this paper.",
            context_text=excerpt,
        )
        search_result = search_papers(rewrite_meta.get("rewritten_query", ""), resolved_limit)
        return {
            "original_query": rewrite_meta.get("original_query", ""),
            "query": rewrite_meta.get("rewritten_query", ""),
            "topics": rewrite_meta.get("topics") or [],
            "reason": rewrite_meta.get("reason", ""),
            "rewrite_model": rewrite_meta.get("model", ""),
            "items": search_result.get("items", []),
            "errors": search_result.get("errors", []),
        }

    def _resolve_section_future(self, future, enabled=True):
        if not enabled:
            return build_section_result("disabled")
        try:
            content_value = future.result()
            if not content_value:
                return build_section_result("empty")
            return build_section_result("success", content=content_value)
        except Exception as exc:
            return build_section_result(
                "failed",
                error=str(exc),
                retryable=is_retryable_llm_error(str(exc)),
            )

    def _finalize_analysis_sections(self, content, sections):
        result = {"char_count": len(content)}
        result["sections"] = sections
        result["summary"] = sections["summary"]["content"]
        result["quotes"] = sections["quotes"]["content"]
        result["mindmap"] = sections["mindmap"]["content"]
        result["mermaid"] = sections["mermaid"]["content"]
        result["evaluation"] = sections["evaluation"]["content"]
        result["research_brief"] = sections["research_brief"]["content"]

        result.update(build_analysis_metadata(sections))

        required_sections = [sections["summary"], sections["quotes"], sections["mindmap"]]
        if all(section["status"] == "failed" for section in required_sections):
            raise RuntimeError(required_sections[0]["error"] or "核心分析项全部失败")
        return result

    def analyze(self, file_path, generate_mermaid=True, generate_evaluation=True, generate_research_brief=True):
        """核心分析流程（已优化为并发执行）"""
        content = DocumentLoader.load(file_path)
        self.document_content = content

        sections = {}
        task_count = 3 + int(generate_mermaid) + int(generate_evaluation) + int(generate_research_brief)
        worker_count = self._get_worker_count(task_count, self.analysis_workers)

        t_start = time.time()

        with concurrent.futures.ThreadPoolExecutor(max_workers=worker_count) as executor:
            future_summary = executor.submit(self.generate_summary, content)
            future_quotes = executor.submit(self.extract_quotes, content)
            future_mindmap = executor.submit(self.generate_mindmap, content)
            future_mermaid = executor.submit(self.generate_mermaid_mindmap, content) if generate_mermaid else None
            future_eval = executor.submit(self.generate_evaluation, content) if generate_evaluation else None
            future_research_brief = executor.submit(self.generate_research_brief, content) if generate_research_brief else None

            sections["summary"] = self._resolve_section_future(future_summary)
            sections["quotes"] = self._resolve_section_future(future_quotes)
            sections["mindmap"] = self._resolve_section_future(future_mindmap)
            sections["mermaid"] = self._resolve_section_future(future_mermaid, enabled=generate_mermaid)
            sections["evaluation"] = self._resolve_section_future(future_eval, enabled=generate_evaluation)
            sections["research_brief"] = self._resolve_section_future(future_research_brief, enabled=generate_research_brief)

        elapsed = time.time() - t_start
        result = self._finalize_analysis_sections(content, sections)
        result["elapsed_seconds"] = round(elapsed, 1)

        logger.info("Analysis completed in %.1fs for %s (%d chars)", elapsed, os.path.basename(file_path), len(content))
        return result

    def analyze_stream(self, file_path, generate_mermaid=True, generate_evaluation=True, generate_research_brief=True):
        content = DocumentLoader.load(file_path)
        self.document_content = content

        sections = {
            "summary": build_section_result("pending"),
            "quotes": build_section_result("pending"),
            "mindmap": build_section_result("pending"),
            "mermaid": build_section_result("disabled") if not generate_mermaid else build_section_result("pending"),
            "evaluation": build_section_result("disabled") if not generate_evaluation else build_section_result("pending"),
            "research_brief": build_section_result("disabled") if not generate_research_brief else build_section_result("pending"),
        }
        task_count = 3 + int(generate_mermaid) + int(generate_evaluation) + int(generate_research_brief)
        worker_count = self._get_worker_count(task_count, self.analysis_workers)
        t_start = time.time()

        with concurrent.futures.ThreadPoolExecutor(max_workers=worker_count) as executor:
            future_map = {
                executor.submit(self.generate_summary, content): ("summary", True),
                executor.submit(self.extract_quotes, content): ("quotes", True),
                executor.submit(self.generate_mindmap, content): ("mindmap", True),
            }
            if generate_mermaid:
                future_map[executor.submit(self.generate_mermaid_mindmap, content)] = ("mermaid", True)
            if generate_evaluation:
                future_map[executor.submit(self.generate_evaluation, content)] = ("evaluation", True)
            if generate_research_brief:
                future_map[executor.submit(self.generate_research_brief, content)] = ("research_brief", True)

            for future in concurrent.futures.as_completed(future_map):
                section_name, enabled = future_map[future]
                section_result = self._resolve_section_future(future, enabled=enabled)
                sections[section_name] = section_result
                yield {
                    "type": "section",
                    "name": section_name,
                    "section": section_result,
                }

        elapsed = time.time() - t_start
        result = self._finalize_analysis_sections(content, sections)
        result["elapsed_seconds"] = round(elapsed, 1)
        logger.info("Streaming analysis completed in %.1fs for %s (%d chars)", elapsed, os.path.basename(file_path), len(content))
        yield {
            "type": "done",
            "result": result,
        }

# ================= 路由控制 =================

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    response = templates.TemplateResponse(request, "index.html")
    response.headers["Cache-Control"] = "no-store"
    return response


@app.get("/logo.ico")
async def logo_ico():
    return FileResponse("logo.ico", media_type="image/x-icon", headers={"Cache-Control": STATIC_ICON_CACHE_CONTROL})


@app.get("/favicon.ico")
async def favicon():
    return FileResponse("logo.ico", media_type="image/x-icon", headers={"Cache-Control": STATIC_ICON_CACHE_CONTROL})


@app.get("/api/health")
async def health_check():
    runtime_folders = {
        "uploads": UPLOAD_FOLDER,
        "output": OUTPUT_FOLDER,
        "context": CONTEXT_FOLDER,
    }
    return JSONResponse(
        content={
            "status": "ok",
            "app": APP_NAME,
            "version": APP_VERSION,
            "timestamp": now_iso(),
            "uptime_seconds": max(0, round(time.time() - APP_STARTED_AT, 3)),
            "folders": {
                name: {
                    "exists": os.path.isdir(path),
                    "writable": os.path.isdir(path) and os.access(path, os.W_OK),
                }
                for name, path in runtime_folders.items()
            },
        },
        headers={"Cache-Control": "no-store"},
    )


@app.post("/api/analyze")
async def analyze(
    file: UploadFile | None = File(None),
    api_key: str = Form(""),
    generate_mermaid: str | None = Form(None),
    generate_evaluation: str | None = Form(None),
    generate_research_brief: str | None = Form(None),
    session_id: str = Form(""),
):
    file_path = None

    if file is None:
        return build_error_response("Please upload a file.", code="missing_file")

    if file.filename == "":
        return build_error_response("Please select a file.", code="missing_filename")
    if not is_allowed_file(file.filename):
        return build_error_response(f"Unsupported file type. Please upload one of: {SUPPORTED_FILE_TYPES_TEXT}", code="unsupported_file_type")

    try:
        cleanup_expired_sessions()
        original_filename = build_safe_upload_filename(file.filename)
        file_path = build_unique_storage_path(UPLOAD_FOLDER, original_filename)

        await save_upload_file(file, file_path, MAX_CONTENT_LENGTH)

        generate_mermaid_bool = parse_bool_value(generate_mermaid, default=True)
        generate_evaluation_bool = parse_bool_value(generate_evaluation, default=True)
        generate_research_brief_bool = parse_bool_value(generate_research_brief, default=True)
        result = analyze_saved_file(
            file_path=file_path,
            original_filename=original_filename,
            api_key=api_key,
            generate_mermaid_bool=generate_mermaid_bool,
            generate_evaluation_bool=generate_evaluation_bool,
            session_id=session_id,
            generate_research_brief_bool=generate_research_brief_bool,
        )
        return JSONResponse(content=result)

    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except Exception as e:
        logger.exception("Document analysis failed")
        return build_error_response(str(e), status_code=500, code="document_analysis_failed")
    finally:
        await close_upload_file_safely(file, "analysis upload file")
        remove_file_safely(file_path, "analyzed upload file")


@app.get("/api/download-paper")
async def download_paper(url: str = "", pdf_url: str = "", title: str = ""):
    try:
        response, file_name, content_type, initial_chunk = stream_remote_paper(title=title, pdf_url=pdf_url, url=url)
        media_type = content_type or mimetypes.guess_type(file_name)[0] or "application/octet-stream"
        quoted_name = urllib.parse.quote(file_name)
        headers = {
            "Content-Disposition": f"attachment; filename=\"{file_name}\"; filename*=UTF-8''{quoted_name}",
            "X-Content-Type-Options": "nosniff",
        }
        return StreamingResponse(iter_remote_file_chunks(response, MAX_CONTENT_LENGTH, initial_chunk), media_type=media_type, headers=headers)
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except urllib.error.HTTPError as exc:
        logger.warning("Paper proxy download failed: %s", exc)
        return build_error_response(describe_remote_http_error(exc.code), code="paper_download_failed")
    except urllib.error.URLError as exc:
        logger.warning("Paper proxy download failed: %s", exc)
        return build_error_response(describe_remote_url_error(getattr(exc, "reason", exc)), code="paper_download_failed")
    except Exception as exc:
        logger.exception("Paper proxy download failed")
        return build_error_response(str(exc), status_code=500, code="paper_download_failed")


@app.post("/api/analyze/stream")
async def analyze_stream(
    file: UploadFile | None = File(None),
    api_key: str = Form(""),
    generate_mermaid: str | None = Form(None),
    generate_evaluation: str | None = Form(None),
    generate_research_brief: str | None = Form(None),
    session_id: str = Form(""),
):
    file_path = None

    if file is None:
        return build_error_response("Please upload a file.", code="missing_file")
    if file.filename == "":
        return build_error_response("Please select a file.", code="missing_filename")
    if not is_allowed_file(file.filename):
        return build_error_response(f"Unsupported file type. Please upload one of: {SUPPORTED_FILE_TYPES_TEXT}", code="unsupported_file_type")

    cleanup_expired_sessions()
    resolved_api_key = resolve_api_key(api_key)
    if not resolved_api_key:
        return build_error_response("API key is required. Provide api_key or set OPENAI_API_KEY.", code="missing_api_key")

    original_filename = build_safe_upload_filename(file.filename)
    file_path = build_unique_storage_path(UPLOAD_FOLDER, original_filename)
    generate_mermaid_bool = parse_bool_value(generate_mermaid, default=True)
    generate_evaluation_bool = parse_bool_value(generate_evaluation, default=True)
    generate_research_brief_bool = parse_bool_value(generate_research_brief, default=True)

    try:
        await save_upload_file(file, file_path, MAX_CONTENT_LENGTH)
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except Exception as exc:
        logger.exception("Streaming analyze upload failed")
        return build_error_response(str(exc), status_code=500, code="streaming_upload_failed")
    finally:
        await close_upload_file_safely(file, "streaming analysis upload file")

    async def event_generator():
        try:
            whisperer = PaperWhisperer(resolved_api_key)
            safe_session_id = build_session_id(session_id)
            yield build_sse_event("start", {"session_id": safe_session_id, "source_filename": original_filename})

            final_result = None
            for event in whisperer.analyze_stream(file_path, generate_mermaid_bool, generate_evaluation_bool, generate_research_brief_bool):
                if event["type"] == "section":
                    yield build_sse_event("section", {"name": event["name"], "section": event["section"]})
                elif event["type"] == "done":
                    final_result = event["result"]

            if final_result is None:
                raise RuntimeError("Analysis stream completed without a final result.")

            final_payload = finalize_analysis_result(
                result=final_result,
                whisperer=whisperer,
                original_filename=original_filename,
                generate_evaluation_bool=generate_evaluation_bool,
                session_id=safe_session_id,
                generate_research_brief_bool=generate_research_brief_bool,
            )
            yield build_sse_event("done", final_payload)
        except Exception as exc:
            logger.exception("Streaming document analysis failed")
            yield build_sse_event("error", build_error_payload(str(exc), code="streaming_analysis_failed"))
        finally:
            remove_file_safely(file_path, "streamed analyzed upload file")

    return StreamingResponse(event_generator(), media_type="text/event-stream", headers=build_sse_headers())


@app.post("/api/import-paper")
async def import_paper(request: Request):
    data, error_response = await parse_json_object_request(request)
    if error_response is not None:
        return error_response

    file_path = None
    try:
        cleanup_expired_sessions()
        title = str(data.get("title") or "").strip()
        url = str(data.get("url") or "").strip()
        pdf_url = str(data.get("pdf_url") or "").strip()
        if not pdf_url and not url:
            return build_error_response("A downloadable paper URL is required.", code="missing_paper_url")

        generate_mermaid_bool = parse_bool_value(data.get("generate_mermaid"), default=True)
        generate_evaluation_bool = parse_bool_value(data.get("generate_evaluation"), default=True)
        generate_research_brief_bool = parse_bool_value(data.get("generate_research_brief"), default=True)
        file_path, original_filename = download_remote_paper(title=title, pdf_url=pdf_url, url=url)
        result = analyze_saved_file(
            file_path=file_path,
            original_filename=original_filename,
            api_key=str(data.get("api_key") or ""),
            generate_mermaid_bool=generate_mermaid_bool,
            generate_evaluation_bool=generate_evaluation_bool,
            session_id=str(data.get("session_id") or ""),
            generate_research_brief_bool=generate_research_brief_bool,
        )
        return JSONResponse(content=result)
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except urllib.error.HTTPError as exc:
        logger.warning("Paper import download failed: %s", exc)
        return build_error_response(describe_remote_http_error(exc.code), code="paper_import_failed")
    except urllib.error.URLError as exc:
        logger.warning("Paper import download failed: %s", exc)
        return build_error_response(describe_remote_url_error(getattr(exc, "reason", exc)), code="paper_import_failed")
    except Exception as exc:
        logger.exception("Paper import failed")
        return build_error_response(str(exc), status_code=500, code="paper_import_failed")
    finally:
        remove_file_safely(file_path, "imported paper analysis file")


@app.post("/api/ask")
async def ask_question(request: Request):
    data, error_response = await parse_json_object_request(request)
    if error_response is not None:
        return error_response

    cleanup_expired_sessions()
    question = str(data.get("question") or "").strip()
    raw_session_id = str(data.get("session_id") or "")
    session_token = str(data.get("session_token") or "")
    answer_mode = normalize_answer_mode(data.get("answer_mode"))

    if not question:
        return build_error_response("Please enter a question.", code="missing_question")

    resolved_api_key = resolve_api_key(data.get("api_key", ""))
    if not resolved_api_key:
        return build_error_response("API key is required. Provide api_key or set OPENAI_API_KEY.", code="missing_api_key")

    try:
        safe_session_id, session_payload = load_validated_session(raw_session_id, session_token, require_token=True)

        whisperer = PaperWhisperer(resolved_api_key)
        whisperer.document_content = get_session_document_content(session_payload)

        t_start = time.time()
        answer = whisperer.answer_question(question, history=session_payload.get("qa_history", []), answer_mode=answer_mode)
        elapsed = time.time() - t_start

        qa_history = session_payload.get("qa_history", [])
        qa_history.append({
            "question": question,
            "answer": answer,
            "answer_mode": answer_mode,
            "timestamp": now_iso(),
        })
        session_payload["qa_history"] = qa_history
        session_payload["document_excerpt"] = build_document_excerpt(get_session_document_content(session_payload))
        write_session_payload(safe_session_id, session_payload)

        logger.info("Q&A completed in %.1fs for session %s", elapsed, safe_session_id)
        return JSONResponse(content={"answer": answer})
    except PermissionError as exc:
        return build_error_response(str(exc), status_code=403, code="invalid_session_token")
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except Exception as e:
        logger.exception("Question answering failed")
        return build_error_response(str(e), status_code=500, code="question_answering_failed")


@app.post("/api/ask/stream")
async def ask_question_stream(request: Request):
    data, error_response = await parse_json_object_request(request)
    if error_response is not None:
        return error_response

    cleanup_expired_sessions()
    question = str(data.get("question") or "").strip()
    raw_session_id = str(data.get("session_id") or "")
    session_token = str(data.get("session_token") or "")
    answer_mode = normalize_answer_mode(data.get("answer_mode"))

    if not question:
        return build_error_response("Please enter a question.", code="missing_question")

    resolved_api_key = resolve_api_key(data.get("api_key", ""))
    if not resolved_api_key:
        return build_error_response("API key is required. Provide api_key or set OPENAI_API_KEY.", code="missing_api_key")

    try:
        safe_session_id, session_payload = load_validated_session(raw_session_id, session_token, require_token=True)
    except PermissionError as exc:
        return build_error_response(str(exc), status_code=403, code="invalid_session_token")
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")

    async def event_generator():
        try:
            whisperer = PaperWhisperer(resolved_api_key)
            whisperer.document_content = get_session_document_content(session_payload)
            yield build_sse_event("start", {"session_id": safe_session_id})

            full_answer_parts = []
            t_start = time.time()
            for chunk in whisperer.stream_answer_question(question, history=session_payload.get("qa_history", []), answer_mode=answer_mode):
                full_answer_parts.append(chunk)
                yield build_sse_event("delta", {"text": chunk})

            answer = "".join(full_answer_parts)
            elapsed = time.time() - t_start

            qa_history = session_payload.get("qa_history", [])
            qa_history.append({
                "question": question,
                "answer": answer,
                "answer_mode": answer_mode,
                "timestamp": now_iso(),
            })
            session_payload["qa_history"] = qa_history
            session_payload["document_excerpt"] = build_document_excerpt(get_session_document_content(session_payload))
            write_session_payload(safe_session_id, session_payload)

            logger.info("Streaming Q&A completed in %.1fs for session %s", elapsed, safe_session_id)
            yield build_sse_event("done", {"answer": answer})
        except Exception as exc:
            logger.exception("Streaming question answering failed")
            yield build_sse_event("error", build_error_payload(str(exc), code="streaming_question_failed"))

    return StreamingResponse(event_generator(), media_type="text/event-stream", headers=build_sse_headers())


@app.post("/api/search-papers")
async def search_papers_api(request: Request):
    data, error_response = await parse_json_object_request(request)
    if error_response is not None:
        return error_response

    cleanup_expired_sessions()
    query = str(data.get("query") or "").strip()
    limit = data.get("limit") or PAPER_SEARCH_RESULT_LIMIT
    raw_session_id = str(data.get("session_id") or "").strip()
    session_token = str(data.get("session_token") or "")
    context_text = str(data.get("context_text") or "")

    if not query:
        return build_error_response("Please enter a search query.", code="missing_search_query")

    try:
        rewrite_meta = {
            "original_query": query,
            "rewritten_query": query,
            "topics": [],
            "reason": "Direct search without AI rewriting.",
            "model": "",
        }
        resolved_api_key = resolve_api_key(data.get("api_key", ""))
        session_payload = None
        safe_session_id = ""

        if raw_session_id:
            safe_session_id, session_payload = load_validated_session(raw_session_id, session_token, require_token=True)

        if PAPER_SEARCH_ENABLE_REWRITE and resolved_api_key:
            whisperer = PaperWhisperer(resolved_api_key)
            rewrite_context = compact_text(context_text, limit=4000)
            if session_payload and not rewrite_context:
                rewrite_context = build_document_excerpt(get_session_document_content(session_payload), limit=4000)
            rewrite_meta = whisperer.rewrite_search_query(query, context_text=rewrite_context)
        elif PAPER_SEARCH_ENABLE_REWRITE:
            rewrite_meta["reason"] = "No API key available; used direct search without AI rewriting."

        result = search_papers(rewrite_meta["rewritten_query"], limit)
        result["original_query"] = rewrite_meta.get("original_query", query)
        result["rewritten_query"] = rewrite_meta.get("rewritten_query", result["query"])
        result["topics"] = rewrite_meta.get("topics", [])
        result["reason"] = rewrite_meta.get("reason", "")
        result["rewrite_model"] = rewrite_meta.get("model", "")

        if session_payload:
            session_payload.setdefault("paper_search", {})
            session_payload["paper_search"]["last_query"] = result["rewritten_query"]
            session_payload["paper_search"]["last_results"] = result["items"]
            write_session_payload(safe_session_id, session_payload)
        return JSONResponse(content=result)
    except PermissionError as exc:
        return build_error_response(str(exc), status_code=403, code="invalid_session_token")
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except Exception as exc:
        logger.exception("Paper search failed")
        return build_error_response(str(exc), status_code=500, code="paper_search_failed")


@app.post("/api/reading-queue")
async def save_reading_queue(request: Request):
    data, error_response = await parse_json_object_request(request)
    if error_response is not None:
        return error_response

    cleanup_expired_sessions()
    raw_session_id = str(data.get("session_id") or "").strip()
    session_token = str(data.get("session_token") or "")

    try:
        safe_session_id, session_payload = load_validated_session(raw_session_id, session_token, require_token=True)
        reading_queue = normalize_paper_collection(data.get("items"), max_items=READING_QUEUE_LIMIT)
        session_payload.setdefault("paper_search", {})
        session_payload["paper_search"]["reading_queue"] = reading_queue
        write_session_payload(safe_session_id, session_payload)
        return JSONResponse(content={"items": reading_queue, "count": len(reading_queue)})
    except PermissionError as exc:
        return build_error_response(str(exc), status_code=403, code="invalid_session_token")
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except Exception as exc:
        logger.exception("Reading queue save failed")
        return build_error_response(str(exc), status_code=500, code="reading_queue_save_failed")


@app.post("/api/recommend-papers")
async def recommend_papers_api(request: Request):
    data, error_response = await parse_json_object_request(request)
    if error_response is not None:
        return error_response

    cleanup_expired_sessions()
    raw_session_id = str(data.get("session_id") or "").strip()
    session_token = str(data.get("session_token") or "")
    if not raw_session_id:
        return build_error_response("session_id is required.", code="missing_session_id")

    limit = data.get("limit") or RECOMMENDATION_RESULT_LIMIT
    resolved_api_key = resolve_api_key(data.get("api_key", ""))
    if not resolved_api_key:
        return build_error_response("API key is required. Provide api_key or set OPENAI_API_KEY.", code="missing_api_key")

    try:
        safe_session_id, session_payload = load_validated_session(raw_session_id, session_token, require_token=True)

        whisperer = PaperWhisperer(resolved_api_key)
        result = whisperer.recommend_papers(get_session_document_content(session_payload), limit=limit)

        session_payload.setdefault("paper_search", {})
        session_payload["paper_search"]["last_recommendation"] = {
            "original_query": result.get("original_query", ""),
            "query": result.get("query", ""),
            "topics": result.get("topics", []),
            "reason": result.get("reason", ""),
            "rewrite_model": result.get("rewrite_model", ""),
            "items": result.get("items", []),
            "errors": result.get("errors", []),
            "generated_at": now_iso(),
        }
        write_session_payload(safe_session_id, session_payload)
        return JSONResponse(content=result)
    except PermissionError as exc:
        return build_error_response(str(exc), status_code=403, code="invalid_session_token")
    except ValueError as exc:
        return build_error_response(str(exc), code="invalid_request")
    except Exception as exc:
        logger.exception("Paper recommendation failed")
        return build_error_response(str(exc), status_code=500, code="paper_recommendation_failed")


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)

    host = os.getenv("FASTAPI_HOST") or os.getenv("FLASK_HOST", "0.0.0.0")
    port = parse_int_env(
        "FASTAPI_PORT",
        default=parse_int_env("FLASK_PORT", default=5000, min_value=1, max_value=65535),
        min_value=1,
        max_value=65535,
    )
    reload_enabled = parse_bool_env("FASTAPI_RELOAD", default=parse_bool_env("FLASK_DEBUG", default=False))

    uvicorn.run("web_app:app", host=host, port=port, reload=reload_enabled)
