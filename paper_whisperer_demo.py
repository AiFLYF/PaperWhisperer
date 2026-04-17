import os
import re
import sys
import time
import threading
import concurrent.futures
from datetime import datetime

from env_loader import load_project_env

load_project_env()

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


SUPPORTED_EXTENSIONS = (".txt", ".pdf", ".docx", ".pptx")
ALLOWED_EXTENSIONS = set(SUPPORTED_EXTENSIONS)
SUPPORTED_FILE_TYPES_TEXT = ", ".join(SUPPORTED_EXTENSIONS)


def parse_int_env(name, default, min_value=1, max_value=32):
    raw = os.getenv(name, "").strip()
    if not raw:
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    return max(min_value, min(max_value, value))


MAX_LLM_CONCURRENCY = parse_int_env("OPENAI_MAX_CONCURRENCY", default=5, min_value=1, max_value=32)
LLM_REQUEST_SEMAPHORE = threading.BoundedSemaphore(MAX_LLM_CONCURRENCY)


def clean_extracted_text(text):
    """Clean extracted text: collapse excessive blank lines and trim whitespace."""
    if not text:
        return ""
    # Collapse 3+ consecutive newlines into 2
    text = re.sub(r'\n{3,}', '\n\n', text)
    # Collapse runs of whitespace (excluding newlines) into single space
    text = re.sub(r'[^\S\n]+', ' ', text)
    # Strip leading/trailing whitespace per line
    lines = [line.strip() for line in text.split('\n')]
    return '\n'.join(lines).strip()


class TextChunker:
    def __init__(self, chunk_size=4000, overlap=200):
        if chunk_size <= 0:
            raise ValueError("chunk_size must be > 0")
        if overlap < 0:
            raise ValueError("overlap must be >= 0")
        if overlap >= chunk_size:
            raise ValueError("overlap must be smaller than chunk_size")
        self.chunk_size = chunk_size
        self.overlap = overlap

    def chunk_text(self, text):
        if len(text) <= self.chunk_size:
            return [text]

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = start + self.chunk_size
            chunk = text[start:end]

            if end < text_length:
                # Try sentence boundaries: Chinese period, English period+space, newline
                best_pos = -1
                for sep in ('。', '. ', '\n'):
                    pos = chunk.rfind(sep)
                    if pos > best_pos:
                        best_pos = pos

                if best_pos > self.chunk_size // 2:
                    chunk = chunk[: best_pos + 1]
                    end = start + best_pos + 1

            chunk = chunk.strip()
            if chunk:
                chunks.append(chunk)
            start = end - self.overlap

        return chunks


class DocumentLoader:
    @staticmethod
    def load_txt(file_path):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def load_pdf(file_path):
        if PdfReader is None:
            raise ValueError("PyPDF2 is not installed. Please install dependencies first.")
        try:
            reader = PdfReader(file_path)
            return "\n\n".join(page.extract_text() or "" for page in reader.pages).strip()
        except Exception as exc:
            raise ValueError(f"Failed to read PDF: {exc}") from exc

    @staticmethod
    def load_docx(file_path):
        if Document is None:
            raise ValueError("python-docx is not installed. Please install dependencies first.")
        try:
            doc = Document(file_path)
            parts = []

            # Extract paragraphs
            for p in doc.paragraphs:
                text = p.text.strip()
                if text:
                    parts.append(text)

            # Extract tables
            for table_index, table in enumerate(doc.tables, start=1):
                rows_text = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append(" | ".join(cells))
                if rows_text:
                    parts.append(f"[Table {table_index}]\n" + "\n".join(rows_text))

            return "\n\n".join(parts).strip()
        except Exception as exc:
            raise ValueError(f"Failed to read DOCX: {exc}") from exc

    @staticmethod
    def load_pptx(file_path):
        if Presentation is None:
            raise ValueError("python-pptx is not installed. Please install dependencies first.")
        try:
            presentation = Presentation(file_path)
            slides_text = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                shape_texts = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        shape_texts.append(text.strip())

                # Extract slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        shape_texts.append(f"[Notes] {notes}")

                if shape_texts:
                    slides_text.append(f"[Slide {slide_index}]\n" + "\n".join(shape_texts))
            return "\n\n".join(slides_text).strip()
        except Exception as exc:
            raise ValueError(f"Failed to read PPTX: {exc}") from exc

    @staticmethod
    def load(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        loaders = {
            ".txt": DocumentLoader.load_txt,
            ".pdf": DocumentLoader.load_pdf,
            ".docx": DocumentLoader.load_docx,
            ".pptx": DocumentLoader.load_pptx,
        }
        loader = loaders.get(ext)
        if not loader:
            raise ValueError(f"Unsupported file type. Please use one of: {SUPPORTED_FILE_TYPES_TEXT}")
        raw_text = loader(file_path)
        return clean_extracted_text(raw_text)


class PaperWhisperer:
    def __init__(self, use_api=True, chunk_size=None, overlap=None):
        self.name = "PaperWhisperer"
        self.version = "0.8.0"
        self.api_key = os.getenv("OPENAI_API_KEY", "").strip()
        self.base_url = os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1").strip()
        self.model = os.getenv("OPENAI_MODEL", "gpt-4o-mini").strip()
        self.request_timeout = parse_int_env("OPENAI_REQUEST_TIMEOUT_SECONDS", default=60, min_value=5, max_value=600)
        self.max_retries = parse_int_env("OPENAI_MAX_RETRIES", default=3, min_value=1, max_value=10)
        self.max_concurrency = MAX_LLM_CONCURRENCY
        self.chunk_workers = min(3, self.max_concurrency)
        self.analysis_workers = min(3, self.max_concurrency)
        self.chunker = TextChunker(
            chunk_size=chunk_size or 4000,
            overlap=overlap or 200,
        )
        self.use_api = bool(use_api and self.api_key and OpenAI is not None)
        self.client = None

        if self.use_api:
            self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)
        elif use_api and OpenAI is None:
            print("[warn] openai package is not installed. Falling back to mock mode.")

    def _call_llm(self, system_prompt, user_prompt, max_retries=None):
        if not self.use_api or not self.client:
            return None

        retries = self.max_retries if max_retries is None else max_retries

        for attempt in range(retries):
            try:
                with LLM_REQUEST_SEMAPHORE:
                    response = self.client.chat.completions.create(
                        model=self.model,
                        messages=[
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_prompt},
                        ],
                        temperature=0.7,
                        max_tokens=4000,
                        timeout=self.request_timeout,
                    )
                return response.choices[0].message.content
            except Exception as exc:
                print(f"[warn] API call failed ({attempt + 1}/{retries}): {exc}")
                if attempt < retries - 1:
                    time.sleep(min(2 * (attempt + 1), 8))

        return None

    def _get_worker_count(self, task_count, configured_workers):
        return max(1, min(task_count, configured_workers))

    def _generate_summary_chunk(self, content):
        system_prompt = "你是专业学术助手。请用中文提炼论文的核心观点，并保持准确简洁。"
        user_prompt = (
            "请阅读以下内容，并输出：\n"
            "1. 3 条核心观点\n"
            "2. 2 条值得引用的片段\n\n"
            f"文档内容：\n{content}\n\n"
            "请按以下格式输出：\n"
            "## 核心观点\n"
            "1. ...\n"
            "2. ...\n"
            "3. ...\n\n"
            "## 引用片段\n"
            "- \"...\"\n"
            "- \"...\""
        )
        return self._call_llm(system_prompt, user_prompt)

    def _merge_summaries(self, summaries):
        if not summaries:
            return None
        if len(summaries) == 1:
            return summaries[0]

        system_prompt = "你是专业学术助手。请将多段论文摘要整合成一份统一、去重、清晰的中文摘要。"
        user_prompt = (
            "以下是同一篇文档不同片段的摘要，请整合输出最终结果。\n\n"
            f"{chr(10).join(summaries)}\n\n"
            "输出格式：\n"
            "## 核心观点\n"
            "1. ...\n"
            "2. ...\n"
            "3. ...\n\n"
            "## 引用片段\n"
            "- \"...\"\n"
            "- \"...\""
        )
        return self._call_llm(system_prompt, user_prompt)

    def generate_summary(self, content):
        print("\n[info] Generating summary...")
        chunks = self.chunker.chunk_text(content)

        if len(chunks) == 1 or not self.use_api:
            result = self._generate_summary_chunk(content) if self.use_api else None
            return result or self._fallback_summary()

        worker_count = self._get_worker_count(len(chunks), self.chunk_workers)
        if worker_count == 1:
            chunk_summaries = [summary for summary in map(self._generate_summary_chunk, chunks) if summary]
        else:
            with concurrent.futures.ThreadPoolExecutor(max_workers=worker_count) as executor:
                chunk_summaries = list(filter(None, executor.map(self._generate_summary_chunk, chunks)))

        merged = self._merge_summaries(chunk_summaries)
        return merged or self._fallback_summary()

    def _fallback_summary(self):
        return (
            "## 核心观点\n"
            "1. 文档整体围绕研究背景、方法设计、实验结果与结论展开。\n"
            "2. 主要价值通常体现在方法创新、实验改进或问题定义上。\n"
            "3. 建议结合原文进一步核对关键公式、数据集与实验设置。\n\n"
            "## 引用片段\n"
            "- \"建议在原文中定位最能代表贡献的结论句进行引用。\"\n"
            "- \"分析结果为辅助阅读结论，正式引用前请再次核对原文。\""
        )

    def extract_quotes(self, content):
        print("\n[info] Extracting quotes...")
        if self.use_api:
            system_prompt = "你是专业学术助手。请从文档中提取适合引用的关键句子，并保持原意。"
            user_prompt = (
                "请从以下内容中提取 3 条值得引用的句子或关键观点：\n\n"
                f"{content[:15000]}\n\n"
                "输出格式：\n"
                "## 引用片段\n"
                "1. \"...\"\n"
                "2. \"...\"\n"
                "3. \"...\""
            )
            result = self._call_llm(system_prompt, user_prompt)
            if result:
                return result

        return (
            "## 引用片段\n"
            "1. \"建议在原文中核对最关键的贡献句。\"\n"
            "2. \"建议优先引用结论、方法概述或实验结果段落。\"\n"
            "3. \"正式写作前请再次核对原文措辞。\""
        )

    def generate_mindmap(self, content):
        print("\n[info] Generating mind map...")
        if self.use_api:
            system_prompt = "你是专业学术助手。请输出中文层级结构，表示文档的主要逻辑结构。"
            user_prompt = (
                "请根据以下文档内容生成简洁的文本思维导图：\n\n"
                f"{content[:10000]}\n\n"
                "输出格式：\n"
                "## 思维导图\n"
                "- 一级主题\n"
                "  - 二级主题"
            )
            result = self._call_llm(system_prompt, user_prompt)
            if result:
                return result

        return (
            "## 思维导图\n"
            "- 研究背景\n"
            "  - 问题定义\n"
            "- 方法设计\n"
            "  - 模型或流程\n"
            "- 实验结果\n"
            "  - 指标表现\n"
            "- 结论\n"
            "  - 价值与局限"
        )

    def save_results(self, file_path, summary, quotes, mindmap, elapsed_seconds=None):
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)

        base_name = os.path.splitext(os.path.basename(file_path))[0]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = os.path.join(output_dir, f"{base_name}_analysis_{timestamp}.md")

        timing_line = f"- Elapsed: {elapsed_seconds:.1f}s\n" if elapsed_seconds is not None else ""

        content = (
            f"# {self.name} Analysis Report\n\n"
            f"> Generated at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"> Source file: {file_path}\n\n"
            "---\n\n"
            f"{summary}\n\n---\n\n{quotes}\n\n---\n\n{mindmap}\n\n---\n\n"
            f"- Version: {self.version}\n"
            f"- API mode: {'enabled' if self.use_api else 'mock'}\n"
            f"- Model: {self.model if self.use_api else 'N/A'}\n"
            f"- Max concurrency: {self.max_concurrency}\n"
            f"{timing_line}"
        )

        with open(output_file, "w", encoding="utf-8") as f:
            f.write(content)

        print(f"\n[info] Saved result to: {output_file}")
        return output_file

    def load_document(self, file_path):
        print(f"[info] Loading document: {file_path}")

        if not os.path.exists(file_path):
            print("[warn] File not found. Using bundled sample content.")
            return self._get_sample_content()

        ext = os.path.splitext(file_path)[1].lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise ValueError(f"Unsupported file type. Please use one of: {SUPPORTED_FILE_TYPES_TEXT}")

        content = DocumentLoader.load(file_path)
        if not content.strip():
            raise ValueError("Document is empty after extraction.")

        print(f"[info] Loaded document successfully ({len(content)} chars)")
        return content

    def _get_sample_content(self):
        return (
            "深度学习在自然语言处理中的应用\n\n"
            "摘要：本文综述了深度学习在 NLP 领域的主要进展，"
            "包括 Transformer 架构、预训练语言模型以及高效推理方法。"
            "研究重点涵盖模型设计、实验结果与未来方向。"
        )

    def run(self, file_path, save_output=True):
        print("=" * 60)
        print(f"{self.name} v{self.version}")
        print("=" * 60)
        print(f"[info] API mode: {'enabled' if self.use_api else 'mock'}")
        if self.use_api:
            print(f"[info] Model: {self.model}")
            print(f"[info] Max concurrency: {self.max_concurrency}")

        t_start = time.time()
        content = self.load_document(file_path)

        # Run analysis tasks concurrently
        if self.use_api and self.analysis_workers > 1:
            worker_count = self._get_worker_count(3, self.analysis_workers)
            with concurrent.futures.ThreadPoolExecutor(max_workers=worker_count) as executor:
                future_summary = executor.submit(self.generate_summary, content)
                future_quotes = executor.submit(self.extract_quotes, content)
                future_mindmap = executor.submit(self.generate_mindmap, content)

                summary = future_summary.result()
                quotes = future_quotes.result()
                mindmap = future_mindmap.result()
        else:
            summary = self.generate_summary(content)
            quotes = self.extract_quotes(content)
            mindmap = self.generate_mindmap(content)

        elapsed = time.time() - t_start

        print("\n" + summary)
        print("\n" + quotes)
        print("\n" + mindmap)
        print(f"\n[info] Analysis complete in {elapsed:.1f}s.")

        if save_output:
            self.save_results(file_path, summary, quotes, mindmap, elapsed_seconds=elapsed)


if __name__ == "__main__":
    file_path = sys.argv[1] if len(sys.argv) > 1 else "sample.txt"
    app = PaperWhisperer(use_api=True)
    app.run(file_path)
